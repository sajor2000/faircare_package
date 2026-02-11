"""
FairCareAI Report Generator

This module generates governance-ready reports in multiple formats:
- PDF: Formal audit report using WeasyPrint
- PPTX: PowerPoint deck for board presentations
- HTML: Standalone interactive dashboard with all 7 sections

Report Structure (per newspec.md):
1. Executive Summary - Results summary
2. Descriptive Statistics - Table 1 cohort summary
3. Overall Performance - TRIPOD+AI metrics
4. Subgroup Performance - Performance by sensitive attribute
5. Fairness Assessment - Fairness metrics analysis
6. Limitations & Flags - Warnings and errors
7. Governance Information - Data for governance team review

Methodology: Van Calster et al. (2025), CHAI RAIC Checkpoint 1.
"""

import asyncio
import html
import re
from concurrent.futures import ThreadPoolExecutor
from dataclasses import dataclass
from datetime import date, datetime
from pathlib import Path
from typing import TYPE_CHECKING, Any

import polars as pl

from faircareai import __version__ as faircareai_version
from faircareai.core.logging import get_logger
from faircareai.visualization.exporters import FigureExportError
from faircareai.visualization.themes import (
    GOVERNANCE_DISCLAIMER_FULL,
    SEMANTIC_COLORS,
    TYPOGRAPHY,
)

if TYPE_CHECKING:
    from faircareai.core.config import MetricDisplayConfig
    from faircareai.core.results import AuditResults
    from faircareai.reports.pptx_options import PptxOptions


logger = get_logger(__name__)

_PLOTLYJS_CDN_URL_FALLBACK = "https://cdn.plot.ly/plotly-2.27.0.min.js"
_PLOTLY_CDN_SCRIPT_RE = re.compile(
    r"<script[^>]+src=[\"']https://cdn\.plot\.ly/plotly-[^\"']+\.min\.js[\"'][^>]*></script>",
    flags=re.IGNORECASE,
)


def _get_plotlyjs_cdn_url() -> str:
    """Return a Plotly.js CDN URL matching the installed plotly.py version."""
    try:
        import plotly.io as pio

        return f"https://cdn.plot.ly/plotly-{pio.plotlyjs_version}.min.js"
    except Exception:
        return _PLOTLYJS_CDN_URL_FALLBACK


def _inject_plotlyjs(html_content: str, *, standalone: bool) -> str:
    """Ensure Plotly.js is available before any embedded figure scripts run.

    Plotly figure fragments are generated with `include_plotlyjs=False` to avoid
    duplicating the library for every chart, so the report must provide Plotly.js
    once at the document level.

    Args:
        html_content: Full HTML document content.
        standalone: If True, inline Plotly.js for offline viewing. If False, link to CDN.
    """
    if "</head>" not in html_content:
        return html_content

    # Remove any existing Plotly CDN tags to avoid duplicates (or stale versions).
    html_content = _PLOTLY_CDN_SCRIPT_RE.sub("", html_content)

    cdn_url = _get_plotlyjs_cdn_url()
    if standalone:
        try:
            from plotly.offline import get_plotlyjs

            plotlyjs = get_plotlyjs()
            script_tag = f'<script type="text/javascript">{plotlyjs}</script>'
        except Exception as e:
            logger.warning(
                "Failed to inline Plotly.js (%s): %s. Falling back to CDN.",
                type(e).__name__,
                str(e),
            )
            script_tag = f'<script src="{cdn_url}"></script>'
    else:
        script_tag = f'<script src="{cdn_url}"></script>'

    return html_content.replace("</head>", f"{script_tag}\n</head>", 1)


def _count_subgroups(results: "AuditResults") -> int:
    """Count demographic subgroups across all sensitive attributes."""
    n_groups = 0
    for _attr_name, metrics in results.subgroup_performance.items():
        if isinstance(metrics, dict):
            groups = metrics.get("groups", metrics)
            n_groups += len(
                [k for k in groups if k not in ("reference", "attribute", "threshold")]
            )
    return n_groups


def _build_audit_trail_rows(
    results: "AuditResults | None",
    summary: "AuditSummary | None",
    report_generated_at: str,
) -> list[tuple[str, str]]:
    """Build audit trail key-value rows for report rendering."""
    audit_id = getattr(results, "audit_id", None) if results else None
    run_timestamp = None
    if results is not None:
        run_timestamp = results.run_timestamp or results.config.report_date
    elif summary is not None:
        run_timestamp = summary.audit_date

    run_timestamp = run_timestamp or date.today().isoformat()

    model_name = (
        results.config.model_name if results is not None else (summary.model_name if summary else "N/A")
    )
    model_version = results.config.model_version if results is not None else ""
    model_label = f"{model_name} v{model_version}" if model_version else model_name

    n_samples = None
    if results is not None:
        n_samples = results.descriptive_stats.get("cohort_overview", {}).get("n_total")
    elif summary is not None:
        n_samples = summary.n_samples

    n_groups = None
    if results is not None:
        n_groups = _count_subgroups(results)
    elif summary is not None:
        n_groups = summary.n_groups

    threshold = None
    if results is not None:
        threshold = results.threshold
    elif summary is not None:
        threshold = summary.threshold

    attributes = "N/A"
    if results is not None:
        attr_names = sorted([k for k in results.fairness_metrics.keys() if isinstance(k, str)])
        attributes = ", ".join(attr_names) if attr_names else "Not specified"

    primary_metric = "Not specified"
    if results is not None and results.config.primary_fairness_metric:
        primary_metric = results.config.primary_fairness_metric.value

    org_name = results.config.organization_name if results is not None else ""

    def _fmt_num(value: Any) -> str:
        if value is None:
            return "N/A"
        try:
            return f"{value:,}"
        except Exception:
            return str(value)

    def _fmt_threshold(value: Any) -> str:
        if value is None:
            return "N/A"
        try:
            return f"{value:.2f}"
        except Exception:
            return str(value)

    rows = [
        ("Audit ID", audit_id or "N/A"),
        ("Audit run timestamp", run_timestamp),
        ("Report generated", report_generated_at),
        ("Model", model_label),
    ]

    if org_name:
        rows.append(("Organization", org_name))

    rows.extend(
        [
            ("Samples", _fmt_num(n_samples)),
            ("Groups", _fmt_num(n_groups)),
            ("Decision threshold", _fmt_threshold(threshold)),
            ("Primary fairness metric", primary_metric),
            ("Sensitive attributes", attributes),
            ("FairCareAI version", faircareai_version),
        ]
    )

    # Reproducibility details
    repro = results.reproducibility if results is not None else {}
    env = repro.get("environment", {}) if isinstance(repro, dict) else {}
    packages = env.get("packages", {}) if isinstance(env, dict) else {}
    python_info = env.get("python", {}) if isinstance(env, dict) else {}
    platform_info = env.get("platform", {}) if isinstance(env, dict) else {}

    def _fmt_packages(keys: list[str]) -> str:
        entries = [f"{k} {packages[k]}" for k in keys if k in packages]
        return ", ".join(entries) if entries else "N/A"

    random_seed = None
    if results is not None:
        random_seed = results.random_seed
    if random_seed is None and isinstance(repro, dict):
        random_seed = repro.get("random_seed")

    if random_seed is not None:
        rows.append(("Random seed", str(random_seed)))

    if python_info:
        python_version = python_info.get("version") or "N/A"
        rows.append(("Python", python_version))

    if platform_info:
        system = platform_info.get("system") or "N/A"
        release = platform_info.get("release") or ""
        rows.append(("OS", f"{system} {release}".strip()))

    key_packages = ["numpy", "polars", "scikit-learn", "statsmodels", "plotly"]
    rows.append(("Key packages", _fmt_packages(key_packages)))

    return rows


def _render_audit_trail_html(
    results: "AuditResults | None",
    summary: "AuditSummary | None",
    report_generated_at: str,
    title: str = "Audit Trail",
) -> str:
    """Render audit trail section as HTML."""
    rows = _build_audit_trail_rows(results, summary, report_generated_at)
    row_html = "\n".join(
        f"<tr><th>{html.escape(label)}</th><td>{html.escape(value)}</td></tr>" for label, value in rows
    )
    return f"""
    <section class="section audit-trail">
        <h2>{html.escape(title)}</h2>
        <table class="audit-table">
            <tbody>
                {row_html}
            </tbody>
        </table>
    </section>
    """


def _is_in_async_context() -> bool:
    """Check if we're running inside an asyncio event loop (e.g., Jupyter notebook)."""
    try:
        asyncio.get_running_loop()
        return True
    except RuntimeError:
        return False


def _run_playwright_pdf_generation(
    html_content: str,
    output_path: Path,
    page_format: str = "Letter",
    margins: dict | None = None,
) -> None:
    """Run Playwright PDF generation, handling async context (Jupyter) safely.

    Args:
        html_content: HTML string to render to PDF.
        output_path: Path for output PDF file.
        page_format: Page format (e.g., "Letter", "A4").
        margins: Page margins dict with top, right, bottom, left keys.
    """
    from playwright.sync_api import sync_playwright

    if margins is None:
        margins = {"top": "0.5in", "right": "0.5in", "bottom": "0.5in", "left": "0.5in"}

    def _generate_pdf() -> None:
        with sync_playwright() as p:
            launch_kwargs: dict[str, Any] = {"headless": True}
            try:
                browser = p.chromium.launch(**launch_kwargs)
            except Exception as e:
                logger.warning(
                    "Playwright Chromium launch failed (%s): %s. Retrying with channel='chrome'.",
                    type(e).__name__,
                    str(e),
                )
                try:
                    browser = p.chromium.launch(channel="chrome", **launch_kwargs)
                except Exception as e2:
                    raise ImportError(
                        "Playwright is installed, but Chromium could not be launched in this environment. "
                        "This is often caused by sandbox/permission restrictions or missing browser binaries. "
                        "Try running outside restricted environments and/or reinstalling browsers with: "
                        "`python -m playwright install chromium`."
                    ) from e2
            page = browser.new_page()

            # Load HTML content with timeout protection (60s for complex reports)
            page.set_content(html_content, wait_until="networkidle", timeout=60000)

            # Ensure Plotly has time to render before PDF capture (prevents blank charts)
            try:
                page.wait_for_function("() => typeof window.Plotly !== 'undefined'", timeout=60000)
                page.wait_for_function(
                    "() => {"
                    "  const charts = Array.from(document.querySelectorAll('.plotly-graph-div'));"
                    "  if (charts.length === 0) return true;"
                    "  return charts.every(c => c.querySelector('svg,canvas'));"
                    "}",
                    timeout=60000,
                )
            except Exception as e:
                logger.warning(
                    "Timed out waiting for charts to render before PDF capture (%s): %s",
                    type(e).__name__,
                    str(e),
                )

            # Generate PDF with print styling
            page.pdf(
                path=str(output_path.resolve()),
                format=page_format,
                margin=margins,
                print_background=True,
            )

            browser.close()

    if _is_in_async_context():
        # Running in Jupyter or other async context - use thread pool
        with ThreadPoolExecutor(max_workers=1) as executor:
            future = executor.submit(_generate_pdf)
            future.result()  # Wait for completion and raise any exceptions
    else:
        # Normal sync context - run directly
        _generate_pdf()


def _validate_output_path(output_path: Path, base_dir: Path | None = None) -> Path:
    """Validate output path is within allowed directory.

    Args:
        output_path: Path to validate
        base_dir: Base directory to restrict writes to. If None, no validation is performed.

    Returns:
        Validated resolved path

    Raises:
        ValueError: If base_dir is provided and output_path is outside base_dir
    """
    resolved = output_path.resolve()

    # Only validate if base_dir is explicitly provided
    if base_dir is not None:
        base = base_dir.resolve()

        # Ensure output path is within base directory
        try:
            resolved.relative_to(base)
        except ValueError:
            raise ValueError(
                f"Security: Output path {resolved} is outside allowed directory {base}. "
                "This could be a path traversal attempt."
            ) from None

    return resolved


@dataclass
class AuditSummary:
    """Container for audit summary data."""

    model_name: str
    audit_date: str
    n_samples: int
    n_groups: int
    threshold: float
    pass_count: int
    warn_count: int
    fail_count: int
    worst_disparity_group: str
    worst_disparity_metric: str
    worst_disparity_value: float
    metrics_df: pl.DataFrame | None = None
    disparities_df: pl.DataFrame | None = None


def generate_pdf_report(
    summary: AuditSummary,
    output_path: str | Path,
    include_charts: bool = True,
    metric_config: "MetricDisplayConfig | None" = None,
    results: "AuditResults | None" = None,
) -> Path:
    """
    Generate a formal PDF audit report.

    Uses Playwright to render HTML with charts to PDF.

    Van Calster et al. (2025) Metric Display:
    -----------------------------------------
    By default, reports show only RECOMMENDED metrics. Pass a MetricDisplayConfig
    with show_optional=True to include OPTIONAL metrics (Brier, O:E ratio, etc.).

    Args:
        summary: AuditSummary with audit results
        output_path: Path for output PDF file
        include_charts: If True, embed charts
        metric_config: MetricDisplayConfig controlling which metrics to display.
            If None, defaults to RECOMMENDED metrics only.
        results: Full AuditResults object for chart generation. If None, charts
            will be limited or unavailable.

    Returns:
        Path to generated PDF file

    Raises:
        ImportError: If Playwright is not installed or chromium browser not available.
            Run: pip install playwright && playwright install chromium
    """
    try:
        from playwright.sync_api import sync_playwright  # noqa: F401
    except ImportError as err:
        raise ImportError(
            "Playwright is required for PDF generation. Install with: "
            "pip install 'faircareai[export]' && playwright install chromium"
        ) from err

    output_path = _validate_output_path(Path(output_path))
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Generate HTML content
    html_content = _generate_report_html(summary, include_charts, results=results)
    html_content = _inject_plotlyjs(html_content, standalone=True)

    # Use Playwright to render HTML to PDF (handles Jupyter/async context)
    _run_playwright_pdf_generation(html_content, output_path)

    return output_path


def generate_pptx_deck(
    summary: AuditSummary,
    output_path: str | Path,
    results: "AuditResults | None" = None,
    include_charts: bool = True,
    pptx_options: "PptxOptions | None" = None,
) -> Path:
    """
    Generate a PowerPoint governance deck.

    Args:
        summary: AuditSummary with audit results
        output_path: Path for output PPTX file

    Returns:
        Path to generated PPTX file
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as err:
        raise ImportError(
            "python-pptx is required for PowerPoint generation. "
            "Install with: pip install 'faircareai[export]'"
        ) from err

    output_path = _validate_output_path(Path(output_path))
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    from faircareai.reports.pptx_options import PptxOptions

    options = pptx_options or PptxOptions()

    if not include_charts:
        options.include_exec_summary_chart = False
        options.include_scorecard_chart = False
        options.include_overall_charts = False
        options.include_subgroup_charts = False
        options.include_vancalster_dashboard = False

    slide_builders: dict[str, callable] = {
        "title": lambda: _add_title_slide(prs, summary, logo_path=options.logo_path),
        "summary": lambda: _add_summary_slide(prs, summary),
        "key_findings": lambda: _add_findings_slide(prs, summary),
        "methodology": lambda: _add_recommendations_slide(prs, summary),
        "exec_summary_chart": lambda: _add_exec_summary_chart_slide(prs, results),
        "scorecard_chart": lambda: _add_scorecard_chart_slide(prs, results),
        "overall_charts": lambda: _add_overall_charts_slide(prs, results),
        "subgroup_charts": lambda: _add_subgroup_charts_slides(prs, results),
        "vancalster_dashboard": lambda: _add_vancalster_slide(prs, results),
    }

    default_order = [
        "title",
        "summary",
        "key_findings",
        "methodology",
        "exec_summary_chart",
        "scorecard_chart",
        "overall_charts",
        "subgroup_charts",
        "vancalster_dashboard",
    ]

    order = options.slide_order or default_order

    for key in order:
        if key == "title" and not options.include_title_slide:
            continue
        if key == "summary" and not options.include_exec_summary:
            continue
        if key == "key_findings" and not options.include_key_findings:
            continue
        if key == "methodology" and not options.include_methodology:
            continue
        if key == "exec_summary_chart" and not options.include_exec_summary_chart:
            continue
        if key == "scorecard_chart" and not options.include_scorecard_chart:
            continue
        if key == "overall_charts" and not options.include_overall_charts:
            continue
        if key == "subgroup_charts" and not options.include_subgroup_charts:
            continue
        if key == "vancalster_dashboard" and not options.include_vancalster_dashboard:
            continue
        if results is None and key in {
            "exec_summary_chart",
            "scorecard_chart",
            "overall_charts",
            "subgroup_charts",
            "vancalster_dashboard",
        }:
            continue
        builder = slide_builders.get(key)
        if builder is None:
            continue
        try:
            builder()
        except ImportError:
            logger.warning(
                "PNG/PPTX chart export requires kaleido. Install with: "
                "pip install 'faircareai[export]'"
            )
        except Exception as exc:
            logger.warning("Failed to add slide '%s': %s", key, exc)

    # Add footer to all slides
    if options.footer_text:
        footer_text = options.footer_text
    else:
        footer_text = f"{summary.model_name} | Audit Date: {summary.audit_date}"
        if results is not None and getattr(results, "audit_id", None):
            footer_text = f"Audit ID: {results.audit_id} | {summary.audit_date}"

    for slide in prs.slides:
        _add_footer(slide, footer_text)

    # Save presentation
    prs.save(str(output_path.resolve()))

    return output_path


def generate_html_report(
    results: "AuditResults",
    output_path: str | Path,
    standalone: bool = True,
    metric_config: "MetricDisplayConfig | None" = None,
) -> Path:
    """
    Generate a comprehensive HTML report with all 7 governance sections.

    If standalone=True, embeds all CSS/JS for offline viewing.

    Van Calster et al. (2025) Metric Display:
    -----------------------------------------
    By default, reports show only RECOMMENDED metrics. Pass a MetricDisplayConfig
    with show_optional=True to include OPTIONAL metrics (Brier, O:E ratio, etc.).

    Args:
        results: AuditResults from FairCareAudit.run()
        output_path: Path for output HTML file
        standalone: If True, embed all assets
        metric_config: MetricDisplayConfig controlling which metrics to display.
            If None, defaults to RECOMMENDED metrics only.

    Returns:
        Path to generated HTML file
    """
    output_path = _validate_output_path(Path(output_path))
    output_path.parent.mkdir(parents=True, exist_ok=True)

    html_content = _generate_full_report_html(results)
    html_content = _inject_plotlyjs(html_content, standalone=standalone)

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(html_content)

    return output_path


def _generate_full_report_html(results: "AuditResults") -> str:
    """Generate comprehensive HTML report with all 7 sections."""

    gov = results.governance_recommendation

    # Determine status
    status = gov.get("status", "REVIEW")
    status_colors = {
        "READY": SEMANTIC_COLORS["pass"],
        "CONDITIONAL": SEMANTIC_COLORS["warn"],
        "REVIEW": SEMANTIC_COLORS["fail"],
    }
    status_color = status_colors.get(status, SEMANTIC_COLORS["fail"])

    report_generated_at = datetime.now().astimezone().isoformat(timespec="seconds")
    audit_run_at = results.run_timestamp or results.config.report_date or date.today().isoformat()

    # Generate sections
    section1_html = _generate_executive_summary_section(results, status, status_color)
    section2_html = _generate_descriptive_section(results)
    section3_html = _generate_performance_section(results)
    section4_html = _generate_subgroup_section(results)
    section5_html = _generate_fairness_section(results)
    section6_html = _generate_flags_section(results)
    section7_html = _generate_governance_section(results)
    audit_trail_html = _render_audit_trail_html(
        results,
        None,
        report_generated_at,
        title="Section 8: Audit Trail",
    )

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>FairCareAI Audit Report: {results.config.model_name}</title>
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');

        :root {{
            --pass-color: {SEMANTIC_COLORS["pass"]};
            --warn-color: {SEMANTIC_COLORS["warn"]};
            --fail-color: {SEMANTIC_COLORS["fail"]};
            --bg-color: #ffffff;
            --text-color: #212529;
            --primary-color: #2c5282;
            --secondary-color: #4a5568;
            --border-color: #e2e8f0;
            --section-bg: #ffffff;
        }}

        * {{ box-sizing: border-box; }}

        /* Scientific Publication Style - Large, Clear, Readable */
        body {{
            font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif;
            font-size: 15px;
            color: var(--text-color);
            background-color: var(--bg-color);
            line-height: 1.45;
            margin: 0;
            padding: 0;
        }}

        .container {{
            max-width: 1000px;
            margin: 0 auto;
            padding: 32px 20px;
        }}

        h1, h2, h3 {{
            font-weight: 600;
            color: var(--primary-color);
            margin-top: 0;
        }}

        /* Publication-style large headers - fixed sizes for HTML readability */
        h1 {{ font-size: 28px; margin-bottom: 10px; }}
        h2 {{ font-size: 20px; margin-top: 32px; border-bottom: 1px solid var(--border-color); padding-bottom: 8px; }}
        h3 {{ font-size: 16px; margin-top: 22px; color: var(--secondary-color); }}

        .header {{
            background: transparent;
            padding: 12px 0 18px 0;
            border-bottom: 2px solid var(--border-color);
            box-shadow: none;
            border-radius: 0;
            margin-bottom: 24px;
        }}

        /* Publication readable metadata */
        .metadata {{ color: #666; font-size: 14px; }}

        .status-badge {{
            display: inline-block;
            padding: 10px 20px;
            border-radius: 6px;
            font-weight: 700;
            font-size: 16px;
            color: white;
            background-color: {status_color};
            margin: 12px 0;
        }}

        .section {{
            background: var(--section-bg);
            padding: 20px;
            border-radius: 6px;
            margin-bottom: 18px;
            border: 1px solid var(--border-color);
            box-shadow: none;
        }}

        .scorecard {{
            display: flex;
            gap: 16px;
            margin: 20px 0;
            flex-wrap: wrap;
        }}

        .scorecard-item {{
            flex: 1;
            min-width: 140px;
            text-align: center;
            padding: 14px;
            border-radius: 6px;
            background: #f9fafb;
        }}

        /* Large scorecard numbers */
        .scorecard-value {{
            font-size: 28px;
            font-weight: 700;
        }}

        .scorecard-label {{
            font-size: 14px;
            color: #666;
            text-transform: uppercase;
            letter-spacing: 0.5px;
        }}

        .pass {{ color: var(--pass-color); }}
        .warn {{ color: var(--warn-color); }}
        .fail {{ color: var(--fail-color); }}

        /* Publication-style readable tables */
        table {{
            width: 100%;
            border-collapse: collapse;
            margin: 16px 0;
            font-size: 14px;
        }}

        th, td {{
            padding: 10px 12px;
            text-align: left;
            border-bottom: 1px solid var(--border-color);
            word-break: break-word;
            white-space: normal;
        }}

        th {{
            background: var(--bg-color);
            font-weight: 600;
            font-size: 15px;
            color: var(--secondary-color);
        }}

        tr:hover {{ background: rgba(0,0,0,0.02); }}

        .flag-item {{
            padding: 12px 16px;
            border-radius: 6px;
            margin: 8px 0;
            border-left: 4px solid;
        }}

        .flag-error {{
            background: rgba(213,94,0,0.1);
            border-color: var(--fail-color);
        }}

        .flag-warning {{
            background: rgba(240,228,66,0.2);
            border-color: var(--warn-color);
        }}

        .governance-block {{
            background: #f7fafc;
            border: 2px solid var(--primary-color);
            padding: 24px;
            border-radius: 8px;
            margin-top: 30px;
        }}



        .footer {{
            margin-top: 40px;
            padding: 20px;
            text-align: center;
            font-size: 14px;
            color: #666;
            border-top: 1px solid var(--border-color);
        }}

        .metric-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(220px, 1fr));
            gap: 20px;
            margin: 20px 0;
        }}

        .metric-card {{
            background: #f9fafb;
            padding: 16px;
            border-radius: 6px;
        }}

        /* Large metric values */
        .metric-value {{
            font-size: 22px;
            font-weight: 700;
            color: var(--primary-color);
        }}

        .metric-label {{
            font-size: 14px;
            color: #666;
        }}

        .chart-placeholder {{
            background: var(--bg-color);
            padding: 40px;
            text-align: center;
            border-radius: 6px;
            color: #666;
            font-size: 16px;
        }}

        /* Responsive chart grid - single column on tablets/mobile */
        @media (max-width: 900px) {{
            .chart-grid, .figure-grid {{
                grid-template-columns: 1fr !important;
            }}
        }}

        /* Figure description boxes - explanatory text above charts */
        .figure-description {{
            background: transparent;
            border-left: 2px solid var(--border-color);
            padding: 8px 12px;
            margin-bottom: 10px;
            font-size: 13px;
            color: #555;
            line-height: 1.4;
            border-radius: 0;
        }}

        .note {{
            border-left: 3px solid var(--primary-color);
            padding: 10px 12px;
            margin: 12px 0;
            font-size: 13px;
            color: #4b5563;
            background: #f9fafb;
            border-radius: 4px;
        }}

        .note-subtle {{
            border-left: 2px solid var(--border-color);
            background: transparent;
            color: #555;
        }}

        .audit-table {{
            width: 100%;
            border-collapse: collapse;
            font-size: 13px;
        }}

        .audit-table th {{
            width: 32%;
            text-align: left;
            padding: 8px 10px;
            background: #f9fafb;
            color: #555;
            font-weight: 600;
            border-bottom: 1px solid var(--border-color);
            word-break: break-word;
            white-space: normal;
        }}

        .audit-table td {{
            padding: 8px 10px;
            border-bottom: 1px solid var(--border-color);
            word-break: break-word;
            white-space: normal;
        }}

        @media print {{
            body {{ background: white; }}
            .section {{ box-shadow: none; border: 1px solid #ddd; }}
            .container {{ max-width: 100%; }}
            .chart-grid, .figure-grid {{ page-break-inside: avoid; }}
        }}
    </style>
</head>
<body>
    <div class="container">
        <header class="header">
            <h1>FairCareAI Audit Report</h1>
            <p class="metadata">
                <strong>Model:</strong> {results.config.model_name} v{results.config.model_version}<br>
                <strong>Report Date:</strong> {results.config.report_date or date.today().isoformat()}<br>
                <strong>Audit Run:</strong> {audit_run_at}<br>
                <strong>Report Generated:</strong> {report_generated_at}<br>
                <strong>Primary Fairness Metric:</strong> {results.config.primary_fairness_metric.value if results.config.primary_fairness_metric else "Not specified"}
            </p>
        </header>

        {section1_html}
        {section2_html}
        {section3_html}
        {section4_html}
        {section5_html}
        {section6_html}
        {section7_html}
        {audit_trail_html}

        <footer class="footer">
            <p>{GOVERNANCE_DISCLAIMER_FULL}</p>
            <p style="font-size: 14px; margin-top: 12px;">
                <strong>Methodology:</strong> Van Calster B, Collins GS, Vickers AJ, et al.
                Evaluation of performance measures in predictive artificial intelligence models
                to support medical decisions: overview and guidance.
                <i>Lancet Digit Health</i> 2025;7(2):e100916.
                DOI: <a href="https://doi.org/10.1016/j.landig.2025.100916" target="_blank" style="color: #2c5282; text-decoration: none;">10.1016/j.landig.2025.100916</a>
            </p>
            <p>Generated by FairCareAI on {report_generated_at}</p>
        </footer>
    </div>
</body>
</html>"""

    return html


def _generate_executive_summary_section(
    results: "AuditResults", status: str, status_color: str
) -> str:
    """Generate Section 1: Executive Summary."""
    gov = results.governance_recommendation
    n_pass = gov.get("n_pass", 0)
    n_warnings = gov.get("n_warnings", 0)
    n_errors = gov.get("n_errors", 0)
    advisory = gov.get("advisory", "")

    status_text = {
        "READY": "Within Threshold",
        "CONDITIONAL": "Near Threshold",
        "REVIEW": "Outside Threshold",
    }

    return f"""
    <section class="section">
        <h2>Section 1: Executive Summary</h2>
        <div class="status-badge">{status_text.get(status, status)}</div>

        <div class="scorecard">
            <div class="scorecard-item">
                <div class="scorecard-value pass">{n_pass}</div>
                <div class="scorecard-label">Pass</div>
            </div>
            <div class="scorecard-item">
                <div class="scorecard-value warn">{n_warnings}</div>
                <div class="scorecard-label">Warnings</div>
            </div>
            <div class="scorecard-item">
                <div class="scorecard-value fail">{n_errors}</div>
                <div class="scorecard-label">Critical</div>
            </div>
        </div>

        <p><strong>Advisory:</strong> {advisory}</p>
        <p style="font-size: 16px; color: #666; font-style: italic;">
            This analysis provides data and metrics for governance committee review. Final deployment decisions are made by the health system governance team through their established review process.
        </p>
    </section>
    """


def _generate_descriptive_section(results: "AuditResults") -> str:
    """Generate Section 2: Descriptive Statistics (Table 1)."""
    desc = results.descriptive_stats
    overview = desc.get("cohort_overview", {})
    pred_dist = desc.get("prediction_distribution", {})

    # Build attribute rows
    attr_rows = ""
    attr_dist = desc.get("attribute_distributions", {})
    outcome_by_attr = desc.get("outcome_by_attribute", {})

    for attr_name, attr_data in attr_dist.items():
        groups = attr_data.get("groups", {})
        outcome_groups = outcome_by_attr.get(attr_name, {}).get("groups", {})
        reference = outcome_by_attr.get(attr_name, {}).get("reference")

        for group_name, group_data in groups.items():
            outcome_data = outcome_groups.get(group_name, {})
            rr = outcome_data.get("rate_ratio")
            rr_str = f"{rr:.2f}" if rr is not None else "ref"
            ref_marker = " (ref)" if group_name == reference else ""

            attr_rows += f"""
            <tr>
                <td>{attr_name}</td>
                <td>{group_name}{ref_marker}</td>
                <td>{group_data.get("n", 0):,}</td>
                <td>{group_data.get("pct_fmt", "N/A")}</td>
                <td>{outcome_data.get("outcome_rate_pct", "N/A")}</td>
                <td>{rr_str}</td>
            </tr>
            """

    return f"""
    <section class="section">
        <h2>Section 2: Descriptive Statistics (Table 1)</h2>

        <h3>Cohort Overview</h3>
        <div class="metric-grid">
            <div class="metric-card">
                <div class="metric-value">{overview.get("n_total", 0):,}</div>
                <div class="metric-label">Total Patients</div>
            </div>
            <div class="metric-card">
                <div class="metric-value">{overview.get("n_positive", 0):,}</div>
                <div class="metric-label">Outcome Positive</div>
            </div>
            <div class="metric-card">
                <div class="metric-value">{overview.get("prevalence_pct", "N/A")}</div>
                <div class="metric-label">Prevalence</div>
            </div>
        </div>

        <h3>Prediction Score Distribution</h3>
        <div class="metric-grid">
            <div class="metric-card">
                <div class="metric-value">{pred_dist.get("mean", 0):.3f}</div>
                <div class="metric-label">Mean Score</div>
            </div>
            <div class="metric-card">
                <div class="metric-value">{pred_dist.get("median", 0):.3f}</div>
                <div class="metric-label">Median Score</div>
            </div>
            <div class="metric-card">
                <div class="metric-value">{pred_dist.get("std", 0):.3f}</div>
                <div class="metric-label">Std Dev</div>
            </div>
        </div>

        <h3>Characteristics by Group</h3>

        <div class="note note-subtle">
            <strong>Note: Sensitive Attribute Stratification</strong>
            <p style="margin: 6px 0 0 0;">
                This analysis stratifies performance by the sensitive attributes defined by the data scientist.
                Each attribute column (e.g., "race_ethnicity", "insurance") is analyzed independently.
                Group values are atomic - for example, "Non-Hispanic White Female Ages 65-74" represents
                ONE demographic group, not separate stratification by race AND age.
            </p>
            <p style="margin: 6px 0 0 0;">
                For cross-attribute analysis (e.g., race × age as 2D matrix), define multiple sensitive attributes separately.
            </p>
        </div>

        <table>
            <thead>
                <tr>
                    <th>Attribute</th>
                    <th>Group</th>
                    <th>N</th>
                    <th>%</th>
                    <th>Outcome Rate</th>
                    <th>Rate Ratio</th>
                </tr>
            </thead>
            <tbody>
                {attr_rows}
            </tbody>
        </table>
    </section>
    """


def _generate_performance_section(results: "AuditResults") -> str:
    """Generate Section 3: Overall Model Performance (TRIPOD+AI)."""
    perf = results.overall_performance
    disc = perf.get("discrimination", {})
    cal = perf.get("calibration", {})
    cls = perf.get("classification_at_threshold", {})

    auroc = disc.get("auroc", 0)
    auroc_ci = disc.get("auroc_ci_fmt", "")
    auprc = disc.get("auprc", 0)
    brier = cal.get("brier_score", 0)
    slope = cal.get("calibration_slope", 1.0)
    threshold = cls.get("threshold", 0.5)

    # Add interpretation guidance
    auroc_interp = (
        "Excellent" if auroc >= 0.8 else "Acceptable" if auroc >= 0.7 else "Below standard"
    )
    brier_interp = (
        "Excellent" if brier < 0.15 else "Acceptable" if brier < 0.25 else "Needs improvement"
    )
    slope_interp = "Well calibrated" if 0.8 <= slope <= 1.2 else "May need recalibration"

    # Generate interactive charts
    charts_html = ""
    roc_html = ""
    prob_dist_html = ""
    dca_html = ""

    try:
        from faircareai.visualization.governance_dashboard import (
            create_governance_overall_figures,
            create_governance_probability_distribution,
            create_governance_roc_curve,
        )

        # EXISTING: 4 gauge figures
        figures = create_governance_overall_figures(results)

        # Extract explanations dict
        explanations = figures.pop("_explanations", {})
        explanation_map = {
            "AUROC": explanations.get("auroc", ""),
            "Calibration": explanations.get("calibration", ""),
            "Brier Score": explanations.get("brier", ""),
            "Classification": explanations.get("classification", ""),
        }

        chart_parts = ['<div class="chart-grid" style="display: grid; grid-template-columns: repeat(2, 1fr); gap: 20px; margin-top: 20px;">']
        for title, fig in figures.items():
            if fig is not None and hasattr(fig, 'to_html'):
                # Add description box above each figure
                desc_text = explanation_map.get(title, "")
                desc_html = f'<div class="figure-description">{html.escape(desc_text)}</div>' if desc_text else ""
                fig_html = fig.to_html(full_html=False, include_plotlyjs=False, div_id=f"chart-{title.replace(' ', '-').lower()}")
                chart_parts.append(f'<div>{desc_html}{fig_html}</div>')
        chart_parts.append('</div>')
        charts_html = ''.join(chart_parts)

        # NEW: Van Calster 2025 figures
        # ROC Curve
        roc_fig = create_governance_roc_curve(results)
        if roc_fig:
            roc_html = f'<div style="margin: 30px 0;">{roc_fig.to_html(full_html=False, include_plotlyjs=False, div_id="chart-roc-curve")}</div>'

        # Probability Distribution
        prob_fig = create_governance_probability_distribution(results)
        if prob_fig:
            prob_dist_html = f'<div style="margin: 30px 0;">{prob_fig.to_html(full_html=False, include_plotlyjs=False, div_id="chart-prob-dist")}</div>'

        # Decision Curve Analysis
        try:
            dca_fig = results.plot_decision_curve()
            if dca_fig:
                dca_html = f'<div style="margin: 30px 0;">{dca_fig.to_html(full_html=False, include_plotlyjs=False, div_id="chart-decision-curve")}</div>'
        except (AttributeError, TypeError) as dca_err:
            logger.warning("Decision curve generation failed: %s", dca_err)
            dca_html = ''

    except (ValueError, TypeError, KeyError) as e:
        logger.warning("Performance chart generation failed: %s", e)
        charts_html = f'<div class="chart-placeholder">Interactive charts could not be generated: {html.escape(str(e))}</div>'
    except ImportError as e:
        logger.error("Chart library not available: %s", e)
        charts_html = '<div class="chart-placeholder">Chart library missing. Install with: pip install \'faircareai[viz]\'</div>'

    return f"""
    <section class="section">
        <h2>Section 3: Overall Model Performance (TRIPOD+AI + Van Calster 2025)</h2>

        <!-- Summary -->
        <div class="note">
            <strong>Summary</strong>
            <p style="margin: 6px 0 0 0;">
                The 4 gauges below provide an at-a-glance summary of model performance across discrimination, calibration, and classification domains.
                Detailed Van Calster 2025 analysis follows.
            </p>
        </div>

        {charts_html}

        <!-- Detailed Van Calster 2025 Analysis -->
        <h3 style="margin-top: 40px; color: #2c5282; border-bottom: 2px solid #2c5282; padding-bottom: 8px;">Detailed Analysis: Van Calster 2025 5-Domain Framework</h3>

        <h4>1. Discrimination: How Well Does the Model Separate Outcomes?</h4>
        <div class="metric-grid">
            <div class="metric-card">
                <div class="metric-value">{auroc:.3f}</div>
                <div class="metric-label">AUROC (0-1 scale) {auroc_ci}</div>
                <div style="margin-top: 8px; font-size: 14px; color: #666;">
                    <strong>{auroc_interp}</strong> (0.5=random, 1.0=perfect)
                </div>
            </div>
            <div class="metric-card">
                <div class="metric-value">{auprc:.3f}</div>
                <div class="metric-label">AUPRC (Precision-Recall)</div>
                <div style="margin-top: 8px; font-size: 14px; color: #666;">
                    Area under precision-recall curve
                </div>
            </div>
        </div>

        {roc_html}

        <p style="color: #666; font-size: 14px; margin: 16px 0;">
            <strong>What this shows:</strong> The ROC curve visualizes how well the model ranks patients from low to high risk.
            A curve hugging the top-left corner means the model correctly identifies high-risk patients without too many false alarms.
            The diagonal line represents random guessing - our model should be well above this line.
        </p>

        {prob_dist_html}

        <p style="color: #666; font-size: 14px; margin: 16px 0;">
            <strong>What this shows:</strong> This chart shows how risk scores differ between patients who develop the outcome (red) versus those who don't (green).
            Good models show clear separation with minimal overlap. Wide gaps indicate strong discrimination.
        </p>

        <h4>2. Calibration: Do Predicted Risks Match Reality?</h4>
        <div class="metric-grid">
            <div class="metric-card">
                <div class="metric-value">{brier:.4f}</div>
                <div class="metric-label">Brier Score (0-0.25, lower=better)</div>
                <div style="margin-top: 8px; font-size: 14px; color: #666;">
                    <strong>{brier_interp}</strong> (&lt;0.15=excellent)
                </div>
            </div>
            <div class="metric-card">
                <div class="metric-value">{slope:.2f}</div>
                <div class="metric-label">Calibration Slope (ideal: 1.00)</div>
                <div style="margin-top: 8px; font-size: 14px; color: #666;">
                    <strong>{slope_interp}</strong> (0.8-1.2=good)
                </div>
            </div>
        </div>

        <h4>3. Classification at Threshold = {threshold:.2f}</h4>
        <p style="color: #666; font-size: 14px; margin-bottom: 16px;">
            At this risk cutoff, here's what happens to patients:
        </p>
        <div class="metric-grid">
            <div class="metric-card">
                <div class="metric-value">{cls.get("sensitivity", 0) * 100:.1f}%</div>
                <div class="metric-label">Sensitivity (TPR)</div>
                <div style="margin-top: 8px; font-size: 14px; color: #666;">
                    % of actual cases correctly identified
                </div>
            </div>
            <div class="metric-card">
                <div class="metric-value">{cls.get("specificity", 0) * 100:.1f}%</div>
                <div class="metric-label">Specificity (TNR)</div>
                <div style="margin-top: 8px; font-size: 14px; color: #666;">
                    % without condition correctly identified
                </div>
            </div>
            <div class="metric-card">
                <div class="metric-value">{cls.get("ppv", 0) * 100:.1f}%</div>
                <div class="metric-label">PPV (Precision)</div>
                <div style="margin-top: 8px; font-size: 14px; color: #666;">
                    When flagged positive, % truly positive
                </div>
            </div>
            <div class="metric-card">
                <div class="metric-value">{cls.get("pct_flagged", 0):.1f}%</div>
                <div class="metric-label">% Flagged High Risk</div>
                <div style="margin-top: 8px; font-size: 14px; color: #666;">
                    Proportion identified for intervention
                </div>
            </div>
        </div>

        <h4>4. Clinical Utility: Does the Model Improve Decisions?</h4>

        <div class="note note-subtle">
            <strong>Decision Curve Analysis</strong>
            <p style="margin: 6px 0;">
                Decision Curve Analysis answers the critical question: "Is using this model better than simple policies?"
                Net benefit measures clinical value per 100 patients at different risk thresholds.
            </p>
            <ul style="margin: 6px 0;">
                <li><strong>Blue line (Model):</strong> Net benefit of using this prediction model</li>
                <li><strong>Gray dashed (Treat All):</strong> Intervene on everyone (wasteful, many false alarms)</li>
                <li><strong>Gray solid (Treat None):</strong> Intervene on no one (misses all cases)</li>
            </ul>
            <p style="margin: 6px 0;">
                <strong>How to interpret:</strong> When the blue model line is above both gray reference lines,
                the model adds clinical value. The range where this occurs shows the "useful threshold range"
                for deployment.
            </p>
        </div>

        {dca_html}
    </section>
    """


def _generate_subgroup_section(results: "AuditResults") -> str:
    """Generate Section 4: Subgroup Performance."""
    subgroup_rows = ""

    for attr_name, attr_data in results.subgroup_performance.items():
        if not isinstance(attr_data, dict):
            continue

        # Extract groups from nested structure
        groups_data = attr_data.get("groups", attr_data)

        for group_name, group_data in groups_data.items():
            # Skip metadata keys
            if group_name in ("attribute", "threshold", "reference", "disparities"):
                continue
            if not isinstance(group_data, dict) or "error" in group_data:
                continue

            auroc = group_data.get("auroc")
            auroc_str = f"{auroc:.3f}" if auroc is not None else "N/A"
            tpr = group_data.get("tpr")
            tpr_str = f"{tpr * 100:.1f}%" if tpr is not None else "N/A"
            fpr = group_data.get("fpr")
            fpr_str = f"{fpr * 100:.1f}%" if fpr is not None else "N/A"
            ref_marker = " (ref)" if group_data.get("is_reference") else ""

            subgroup_rows += f"""
            <tr>
                <td>{attr_name}</td>
                <td>{group_name}{ref_marker}</td>
                <td>{group_data.get("n", 0):,}</td>
                <td>{auroc_str}</td>
                <td>{tpr_str}</td>
                <td>{fpr_str}</td>
            </tr>
            """

    # Generate interactive subgroup charts
    charts_html = ""
    try:
        from faircareai.visualization.governance_dashboard import (
            create_governance_subgroup_figures,
        )

        all_figures = create_governance_subgroup_figures(results)
        # Chart explanations (shown as HTML below each chart for better spacing)
        CHART_EXPLANATIONS = {
            "Model Accuracy (AUROC) by Demographic Group": (
                "Does the model perform equally well across all groups? "
                "All bars should be similar height (difference <0.05 is ideal). "
                "Lower bars = less accurate for that group."
            ),
            "Sensitivity: % of Actual Cases Detected by Group": (
                "Of patients who develop the outcome, what % does the model correctly identify? "
                "Fairness goal: Differences between groups should be <10 percentage points."
            ),
            "False Alarms: % Incorrectly Flagged by Group": (
                "Of patients without the outcome, what % are incorrectly flagged? "
                "Lower is better. Higher FPR = more unnecessary interventions for that group."
            ),
            "Intervention Rate: % Flagged as High-Risk by Group": (
                "What % of each group is flagged as high-risk? "
                "Large differences may indicate disparate treatment."
            ),
        }

        chart_parts = []
        for attr_name, figures in all_figures.items():
            chart_parts.append(f'<h3 style="margin-top: 30px; color: #2c5282;">{attr_name.replace("_", " ").title()}</h3>')
            chart_parts.append('<div class="subgroup-charts" style="display: flex; flex-direction: column; gap: 40px; margin-top: 20px;">')
            for title, fig in figures.items():
                if fig is not None:
                    fig_html = fig.to_html(full_html=False, include_plotlyjs=False, div_id=f"chart-{attr_name}-{title.replace(' ', '-').lower()}")
                    explanation = CHART_EXPLANATIONS.get(title, "")
                    explanation_html = f'<p style="color: #666; font-size: 13px; margin-top: 8px; padding: 0 20px;">{explanation}</p>' if explanation else ""
                    chart_parts.append(f'<div style="margin-bottom: 20px;">{fig_html}{explanation_html}</div>')
            chart_parts.append('</div>')
        charts_html = ''.join(chart_parts)
    except (ValueError, TypeError, KeyError) as e:
        logger.warning("Subgroup chart generation failed: %s", e)
        charts_html = f'<div class="chart-placeholder">Interactive charts could not be generated: {html.escape(str(e))}</div>'
    except ImportError as e:
        logger.error("Chart library not available: %s", e)
        charts_html = '<div class="chart-placeholder">Chart library missing. Install with: pip install \'faircareai[viz]\'</div>'

    return f"""
    <section class="section">
        <h2>Section 4: Subgroup Performance</h2>

        <p style="color: #666; font-size: 16px; margin-bottom: 20px;">
            <strong>What to look for:</strong> Performance should be similar across all demographic groups.
            Large differences in AUROC (&gt;0.05) or TPR/FPR (&gt;10 percentage points) may indicate fairness concerns.
        </p>

        <table>
            <thead>
                <tr>
                    <th>Attribute</th>
                    <th>Group</th>
                    <th>Sample Size</th>
                    <th>AUROC<br><span style="font-weight: normal; font-size: 12px;">(accuracy)</span></th>
                    <th>TPR<br><span style="font-weight: normal; font-size: 12px;">(sensitivity)</span></th>
                    <th>FPR<br><span style="font-weight: normal; font-size: 12px;">(false alarms)</span></th>
                </tr>
            </thead>
            <tbody>
                {subgroup_rows}
            </tbody>
        </table>

        <div style="margin-top: 20px; padding: 16px; background: #f8f9fa; border-left: 4px solid #0072B2; border-radius: 4px;">
            <h4 style="margin-top: 0; color: #0072B2;">Interpreting These Metrics:</h4>
            <ul style="margin-bottom: 0;">
                <li><strong>AUROC:</strong> Model's ability to rank patients (0.7+ acceptable, 0.8+ strong)</li>
                <li><strong>TPR (Sensitivity):</strong> % of actual cases caught by the model (higher is better)</li>
                <li><strong>FPR:</strong> % incorrectly flagged (lower is better, means fewer false alarms)</li>
                <li><strong>(ref):</strong> Reference group used for fairness comparisons</li>
            </ul>
        </div>

        {charts_html}
    </section>
    """


def _generate_fairness_section(results: "AuditResults") -> str:
    """Generate Section 5: Fairness Assessment with metric-specific content."""
    from faircareai.core.config import FairnessMetric

    config = results.config
    metric = config.primary_fairness_metric
    justification = config.fairness_justification or "Not provided"

    # Metric-specific descriptions and what to look for
    metric_info = {
        FairnessMetric.DEMOGRAPHIC_PARITY: {
            "name": "Demographic Parity",
            "description": "Equal selection rates across groups regardless of true outcomes.",
            "what_to_look_for": "Selection rate differences should be small. Large differences mean some groups are selected more/less often.",
            "key_metric": "selection_rate_diff",
            "threshold_note": "Differences < 0.10 (10%) are typically acceptable.",
        },
        FairnessMetric.EQUALIZED_ODDS: {
            "name": "Equalized Odds",
            "description": "Equal true positive rates AND false positive rates across groups.",
            "what_to_look_for": "Both TPR and FPR differences should be small. This ensures equal benefit AND equal burden across groups.",
            "key_metric": "equalized_odds",
            "threshold_note": "Max(TPR diff, FPR diff) < 0.10 is typically acceptable.",
        },
        FairnessMetric.EQUAL_OPPORTUNITY: {
            "name": "Equal Opportunity",
            "description": "Equal true positive rates across groups (focuses on benefit, not burden).",
            "what_to_look_for": "TPR differences should be small. This ensures all groups with the condition are equally likely to be identified.",
            "key_metric": "equal_opportunity",
            "threshold_note": "TPR differences < 0.10 are typically acceptable.",
        },
        FairnessMetric.PREDICTIVE_PARITY: {
            "name": "Predictive Parity",
            "description": "Equal positive predictive value (PPV) across groups.",
            "what_to_look_for": "PPV differences should be small. A positive prediction should mean the same thing for all groups.",
            "key_metric": "ppv_diff",
            "threshold_note": "PPV differences < 0.10 are typically acceptable.",
        },
        FairnessMetric.CALIBRATION: {
            "name": "Calibration",
            "description": "Predicted probabilities match actual outcomes equally across groups.",
            "what_to_look_for": "Calibration error differences should be small. A 30% prediction should mean 30% risk for all groups.",
            "key_metric": "calibration_diff",
            "threshold_note": "Calibration differences < 0.05 are typically acceptable.",
        },
    }

    # Get info for selected metric
    selected_info = metric_info.get(metric, {
        "name": "Not Specified",
        "description": "No primary fairness metric selected.",
        "what_to_look_for": "Review all metrics below.",
        "key_metric": None,
        "threshold_note": "Differences < 0.10 are typically acceptable.",
    })

    # Build table rows with all metrics, highlighting the primary one
    fairness_rows = ""
    for attr_name, attr_data in results.fairness_metrics.items():
        if not isinstance(attr_data, dict):
            continue

        summary = attr_data.get("summary", {})

        # Demographic Parity (selection rate)
        dp = summary.get("demographic_parity", {})
        dp_diff = dp.get("worst_diff", 0) if dp else 0
        dp_pass = dp.get("within_threshold", True) if dp else True

        # Equal Opportunity (TPR)
        eo = summary.get("equal_opportunity", {})
        eo_diff = eo.get("worst_diff", 0) if eo else 0
        eo_pass = eo.get("within_threshold", True) if eo else True

        # Equalized Odds (TPR + FPR)
        eq = summary.get("equalized_odds", {})
        eq_diff = eq.get("worst_diff", 0) if eq else 0
        eq_pass = eq.get("within_threshold", True) if eq else True

        # Predictive Parity (PPV)
        pp = summary.get("predictive_parity", {})
        pp_diff = pp.get("worst_diff", 0) if pp else 0
        pp_pass = pp.get("within_threshold", True) if pp else True

        # Calibration
        cal = summary.get("calibration", {})
        cal_diff = cal.get("worst_diff", 0) if cal else 0
        cal_pass = cal.get("within_threshold", True) if cal else True

        # Helper to format cell with highlighting for primary metric
        def format_cell(value: float, passed: bool, is_primary: bool) -> str:
            status = "PASS" if passed else "FLAG"
            status_class = "pass" if passed else "fail"
            highlight = ' style="background: #e8f4f8; font-weight: bold;"' if is_primary else ""
            return f'<td{highlight}>{abs(value):.3f}</td><td class="{status_class}"{highlight}>{status}</td>'

        # Determine which metric is primary for this row
        is_dp_primary = metric == FairnessMetric.DEMOGRAPHIC_PARITY
        is_eo_primary = metric == FairnessMetric.EQUAL_OPPORTUNITY
        is_eq_primary = metric == FairnessMetric.EQUALIZED_ODDS
        is_pp_primary = metric == FairnessMetric.PREDICTIVE_PARITY
        is_cal_primary = metric == FairnessMetric.CALIBRATION

        fairness_rows += f'''
        <tr>
            <td>{attr_name}</td>
            {format_cell(dp_diff, dp_pass, is_dp_primary)}
            {format_cell(eo_diff, eo_pass, is_eo_primary)}
            {format_cell(eq_diff, eq_pass, is_eq_primary)}
            {format_cell(pp_diff, pp_pass, is_pp_primary)}
            {format_cell(cal_diff, cal_pass, is_cal_primary)}
        </tr>
        '''

    # Primary metric badge color
    metric_color = "#0072B2" if metric else "#666"

    return f"""
    <section class="section">
        <h2>Section 5: Fairness Assessment</h2>

        <div class="note" style="border-left-color: {metric_color};">
            <strong>Primary Fairness Metric: {selected_info["name"]}</strong>
            <p style="margin: 6px 0 0 0;"><strong>Definition:</strong> {selected_info["description"]}</p>
            <p style="margin: 6px 0 0 0;"><strong>What to look for:</strong> {selected_info["what_to_look_for"]}</p>
            <p style="margin: 6px 0 0 0;"><strong>Threshold:</strong> {selected_info["threshold_note"]}</p>
            <p style="margin: 6px 0 0 0;"><strong>Justification:</strong> {justification}</p>
        </div>

        <h3>All Fairness Metrics by Attribute</h3>
        <p style="color: #666; font-size: 14px; margin-bottom: 16px;">
            Your selected metric is <strong>highlighted in blue</strong>. Other metrics shown for completeness.
        </p>

        <div style="overflow-x: auto;">
        <table style="font-size: 14px;">
            <thead>
                <tr>
                    <th>Attribute</th>
                    <th colspan="2">Demographic Parity<br><span style="font-weight: normal; font-size: 12px;">Selection Rate Diff</span></th>
                    <th colspan="2">Equal Opportunity<br><span style="font-weight: normal; font-size: 12px;">TPR Diff</span></th>
                    <th colspan="2">Equalized Odds<br><span style="font-weight: normal; font-size: 12px;">Max(TPR, FPR) Diff</span></th>
                    <th colspan="2">Predictive Parity<br><span style="font-weight: normal; font-size: 12px;">PPV Diff</span></th>
                    <th colspan="2">Calibration<br><span style="font-weight: normal; font-size: 12px;">Cal Error Diff</span></th>
                </tr>
            </thead>
            <tbody>
                {fairness_rows}
            </tbody>
        </table>
        </div>

        <div class="note note-subtle" style="margin-top: 16px;">
            <strong>Why your metric choice matters:</strong>
            <p style="margin: 6px 0;">The <strong>impossibility theorem</strong> proves that when base rates differ between groups,
            no model can satisfy all fairness criteria simultaneously. Your choice reflects your values:</p>
            <ul style="margin: 6px 0 0 0;">
                <li><strong>Demographic Parity:</strong> Prioritizes equal selection rates (good for resource allocation)</li>
                <li><strong>Equal Opportunity:</strong> Prioritizes equal detection of true cases (good for screening)</li>
                <li><strong>Equalized Odds:</strong> Balances detection AND false alarms (good for interventions)</li>
                <li><strong>Predictive Parity:</strong> Prioritizes equal meaning of positive predictions</li>
                <li><strong>Calibration:</strong> Prioritizes accurate risk communication across groups</li>
            </ul>
        </div>
    </section>
    """


def _generate_flags_section(results: "AuditResults") -> str:
    """Generate Section 6: Flags and Warnings."""
    flag_parts = []

    for flag in results.flags:
        severity = flag.get("severity", "warning")
        flag_class = "flag-error" if severity == "error" else "flag-warning"
        message = flag.get("message", "")
        details = flag.get("details", "")

        flag_parts.append(f"""
        <div class="flag-item {flag_class}">
            <strong>{severity.upper()}:</strong> {html.escape(message)}
            {f"<br><small>{html.escape(details)}</small>" if details else ""}
        </div>
        """)

    flags_html = ''.join(flag_parts) if flag_parts else '<p style="color: green;">No flags or warnings raised.</p>'

    return f"""
    <section class="section">
        <h2>Section 6: Flags and Warnings</h2>
        {flags_html}
    </section>
    """


def _generate_governance_section(results: "AuditResults") -> str:
    """Generate Section 7: Governance Decision Block."""
    config = results.config
    gov = results.governance_recommendation

    return f"""
    <section class="section governance-block">
        <h2>Section 7: Governance Decision Support</h2>

        <p><strong>Model:</strong> {config.model_name} v{config.model_version}</p>
        <p><strong>Intended Use:</strong> {config.intended_use or "Not specified"}</p>
        <p><strong>Intended Population:</strong> {config.intended_population or "Not specified"}</p>
        <p><strong>System Status:</strong> {gov.get("status", "REVIEW")}</p>

        <div class="governance-note" style="background-color: #f8f9fa; border-left: 4px solid #0072B2; padding: 16px; margin-top: 20px;">
            <h3>Governance Process Note</h3>
            <p>
                This package provides essential fairness and performance data to support the
                <strong>Health System Governance Team's</strong> adjudication process.
            </p>
            <p>
                The metrics and advisory status above are calculated outputs based on configured thresholds.
                <strong>Final deployment decisions are made by the Governance Team outside of this technical package.</strong>
            </p>
            <p>
                Please refer to your organization's standard operating procedures for voting, adjudication,
                and final sign-off documentation.
            </p>
        </div>

        <p style="font-size: 16px; color: #666; margin-top: 24px; font-style: italic;">
            {GOVERNANCE_DISCLAIMER_FULL}
        </p>
    </section>
    """


def _generate_report_html(
    summary: AuditSummary,
    include_charts: bool = True,
    results: "AuditResults | None" = None,
) -> str:
    """Generate the HTML content for the report.

    Args:
        summary: AuditSummary with basic audit info.
        include_charts: If True, generate charts.
        results: Full AuditResults for chart generation. If None, charts will be limited.
    """

    from faircareai.visualization.tables import create_plain_language_summary

    # Determine overall status based on threshold results
    if summary.fail_count > 0:
        overall_status = "Outside Threshold"
        status_color = SEMANTIC_COLORS["fail"]
    elif summary.warn_count > 0:
        overall_status = "Near Threshold"
        status_color = SEMANTIC_COLORS["warn"]
    else:
        overall_status = "Within Threshold"
        status_color = SEMANTIC_COLORS["pass"]

    # Generate plain language summary
    plain_summary = create_plain_language_summary(
        summary.pass_count,
        summary.warn_count,
        summary.fail_count,
        summary.worst_disparity_group,
        summary.worst_disparity_metric,
        summary.worst_disparity_value,
    )

    report_generated_at = datetime.now().astimezone().isoformat(timespec="seconds")
    audit_run_at = (
        results.run_timestamp if results is not None else None
    ) or summary.audit_date or date.today().isoformat()

    # Generate charts if requested
    charts_html = ""
    if include_charts:
        if results is not None:
            # Use full AuditResults for comprehensive charts
            try:
                overall_html = _render_governance_overall_figures(results)
                subgroup_html = _render_governance_subgroup_figures(results)
                charts_html = f"""
                <div class="charts-section">
                    <h3>Overall Performance</h3>
                    {overall_html}
                    <h3>Subgroup Performance</h3>
                    {subgroup_html}
                </div>
                """
            except (ValueError, TypeError, KeyError) as e:
                logger.warning("Chart generation failed: %s", e)
                charts_html = (
                    f'<p class="chart-placeholder">Charts could not be generated: {html.escape(str(e))}</p>'
                )
            except ImportError as e:
                logger.error("Chart library not available: %s", e)
                charts_html = '<p class="chart-placeholder">Chart library missing. Install with: pip install \'faircareai[viz]\'</p>'
        elif summary.metrics_df is not None and len(summary.metrics_df) > 0:
            # Fall back to forest plot from metrics_df
            try:
                from faircareai.visualization.altair_plots import create_forest_plot_static

                chart = create_forest_plot_static(summary.metrics_df, metric="tpr")
                charts_html = f'<div class="chart-container">{chart.to_html()}</div>'
            except (ValueError, TypeError, KeyError) as e:
                logger.warning("Forest plot generation failed: %s", e)
                charts_html = '<p class="chart-placeholder">Charts could not be generated.</p>'
            except ImportError as e:
                logger.error("Chart library not available: %s", e)
                charts_html = '<p class="chart-placeholder">Chart library missing. Install with: pip install \'faircareai[viz]\'</p>'
        else:
            charts_html = '<p class="chart-placeholder">No chart data available.</p>'

    audit_trail_html = _render_audit_trail_html(
        results,
        summary,
        report_generated_at,
        title="Audit Trail",
    )

    html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>FairCare Equity Audit Report: {summary.model_name}</title>
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Merriweather:wght@300;400;700&family=Inter:wght@400;500;600&display=swap');

        :root {{
            --pass-color: {SEMANTIC_COLORS["pass"]};
            --warn-color: {SEMANTIC_COLORS["warn"]};
            --fail-color: {SEMANTIC_COLORS["fail"]};
            --bg-color: #ffffff;
            --text-color: {SEMANTIC_COLORS["text"]};
        }}

        * {{
            box-sizing: border-box;
        }}

        body {{
            font-family: {TYPOGRAPHY["data_font"]};
            font-size: 15px;
            color: var(--text-color);
            background-color: var(--bg-color);
            line-height: 1.45;
            max-width: 900px;
            margin: 0 auto;
            padding: 32px 20px;
        }}

        h1, h2, h3 {{
            font-family: {TYPOGRAPHY["heading_font"]};
            font-weight: {TYPOGRAPHY["heading_weight"]};
            color: var(--text-color);
        }}

        h1 {{
            font-size: 32px;
            margin-bottom: 8px;
        }}

        h2 {{
            font-size: 24px;
            margin-top: 40px;
            border-bottom: 2px solid var(--text-color);
            padding-bottom: 8px;
        }}

        .header {{
            border-bottom: 3px solid var(--text-color);
            padding-bottom: 20px;
            margin-bottom: 30px;
        }}

        .metadata {{
            color: #666;
            font-size: 14px;
        }}

        .status-badge {{
            display: inline-block;
            padding: 8px 20px;
            border-radius: 4px;
            font-weight: bold;
            font-size: 18px;
            color: white;
            background-color: {status_color};
            margin: 16px 0;
        }}

        .scorecard {{
            display: flex;
            gap: 20px;
            margin: 24px 0;
        }}

        .scorecard-item {{
            flex: 1;
            text-align: center;
            padding: 20px;
            border-radius: 8px;
            background: white;
            box-shadow: 0 1px 3px rgba(0,0,0,0.1);
        }}

        .scorecard-value {{
            font-size: 36px;
            font-weight: bold;
        }}

        .scorecard-label {{
            font-size: 14px;
            color: #666;
            text-transform: uppercase;
        }}

        .pass {{ color: var(--pass-color); }}
        .warn {{ color: var(--warn-color); }}
        .fail {{ color: var(--fail-color); }}

        .summary-section {{
            background: white;
            padding: 24px;
            border-radius: 8px;
            margin: 24px 0;
            box-shadow: 0 1px 3px rgba(0,0,0,0.1);
        }}

        .chart-container {{
            margin: 24px 0;
        }}

        .footer {{
            margin-top: 60px;
            padding-top: 20px;
            border-top: 1px solid #ddd;
            font-size: 12px;
            color: #666;
        }}

        .audit-table {{
            width: 100%;
            border-collapse: collapse;
            font-size: 13px;
            margin-top: 12px;
        }}

        .audit-table th {{
            width: 32%;
            text-align: left;
            padding: 8px 10px;
            background: #f5f5f5;
            color: #555;
            font-weight: 600;
            border-bottom: 1px solid #ddd;
            word-break: break-word;
            white-space: normal;
        }}

        .audit-table td {{
            padding: 8px 10px;
            border-bottom: 1px solid #ddd;
            word-break: break-word;
            white-space: normal;
        }}

        @media print {{
            body {{
                max-width: 100%;
                padding: 20px;
            }}
        }}

        .charts-section {{
            margin: 30px 0;
        }}

        .charts-section h3 {{
            font-size: 20px;
            margin-top: 30px;
            margin-bottom: 15px;
            color: #2c5282;
        }}
    </style>
</head>
<body>
    <header class="header">
        <h1>Equity Audit Report</h1>
        <p class="metadata">
            Model: <strong>{summary.model_name}</strong><br>
            Audit Date: {summary.audit_date}<br>
            Audit Run: {audit_run_at}<br>
            Report Generated: {report_generated_at}<br>
            Samples: {summary.n_samples:,} | Groups: {summary.n_groups} | Threshold: {summary.threshold:.0%}
        </p>
    </header>

    <section>
        <div class="status-badge">{overall_status}</div>

        <div class="scorecard">
            <div class="scorecard-item">
                <div class="scorecard-value pass">{summary.pass_count}</div>
                <div class="scorecard-label">Pass</div>
            </div>
            <div class="scorecard-item">
                <div class="scorecard-value warn">{summary.warn_count}</div>
                <div class="scorecard-label">Review</div>
            </div>
            <div class="scorecard-item">
                <div class="scorecard-value fail">{summary.fail_count}</div>
                <div class="scorecard-label">Flag</div>
            </div>
        </div>
        <p style="font-size: 16px; color: #666; margin-top: 8px;">
            <em>Data provided for governance team review — final decisions made by health system governance committee</em>
        </p>
    </section>

    <section class="summary-section">
        {plain_summary}
    </section>

    <h2>Detailed Analysis</h2>
    {charts_html}
    {audit_trail_html}

    <footer class="footer">
        <p>
            Generated by FairCareAI |
            Report generated on {report_generated_at}
        </p>
    </footer>
</body>
</html>"""

    return html_content


def _get_print_css() -> str:
    """Get CSS for print/PDF output."""
    return """
    @page {
        size: letter;
        margin: 1in;

        @top-right {
            content: "FairCareAI Audit Report";
            font-size: 12pt;
            color: #666;
        }

        @bottom-center {
            content: counter(page) " of " counter(pages);
            font-size: 12pt;
            color: #666;
        }
    }

    @page :first {
        @top-right { content: none; }
    }

    body {
        font-size: 12pt;
    }

    h1 { page-break-after: avoid; }
    h2 { page-break-after: avoid; }

    .scorecard {
        page-break-inside: avoid;
    }

    .chart-container {
        page-break-inside: avoid;
    }
    """


def _add_title_slide(prs: Any, summary: AuditSummary, logo_path: str | Path | None = None) -> None:
    """Add title slide to presentation with publication-ready typography."""
    from contextlib import suppress

    from pptx.util import Inches, Pt

    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title - large, clear (44pt from TYPOGRAPHY)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Equity Audit Report"
    p.font.size = Pt(TYPOGRAPHY["ppt_title_size"])  # 44pt
    p.font.bold = True

    # Subtitle - readable (32pt)
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = summary.model_name
    p.font.size = Pt(TYPOGRAPHY["ppt_subtitle_size"])  # 32pt

    # Date - readable body text (24pt)
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Audit Date: {summary.audit_date}"
    p.font.size = Pt(TYPOGRAPHY["ppt_body_size"])  # 24pt

    if logo_path:
        with suppress(Exception):
            _add_logo(slide, logo_path)


def _add_logo(slide: Any, logo_path: str | Path) -> None:
    """Add a logo image to the top-right corner."""
    from pptx.util import Inches

    logo = Path(logo_path)
    if not logo.exists():
        return
    slide.shapes.add_picture(str(logo), Inches(10.8), Inches(0.3), width=Inches(2.0))


def _add_footer(slide: Any, footer_text: str) -> None:
    """Add a small footer to a slide."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = footer_text
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.RIGHT


def _add_slide_title(slide: Any, title: str) -> None:
    """Add a slide title with consistent typography."""
    from pptx.util import Inches, Pt

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(TYPOGRAPHY["headline_size"])
    p.font.bold = True


def _add_summary_slide(prs: Any, summary: AuditSummary) -> None:
    """Add executive summary slide with publication-ready typography."""
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title - publication style (36pt from subheading)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Executive Summary"
    p.font.size = Pt(TYPOGRAPHY["headline_size"])  # 36pt
    p.font.bold = True

    # Status based on threshold results
    if summary.fail_count > 0:
        status = "Outside Threshold"
        color = RGBColor(0xD5, 0x5E, 0x00)  # Vermillion
    elif summary.warn_count > 0:
        status = "Near Threshold"
        color = RGBColor(0xF0, 0xE4, 0x42)  # Yellow
    else:
        status = "Within Threshold"
        color = RGBColor(0x00, 0x9E, 0x73)  # Green

    status_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(0.7))
    tf = status_box.text_frame
    p = tf.paragraphs[0]
    p.text = status
    p.font.size = Pt(TYPOGRAPHY["ppt_subtitle_size"])  # 32pt - prominent status
    p.font.bold = True
    p.font.color.rgb = color

    # Metrics - using advisory terminology (large readable)
    metrics_text = (
        f"PASS: {summary.pass_count}    REVIEW: {summary.warn_count}    FLAG: {summary.fail_count}"
    )
    metrics_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(0.5))
    tf = metrics_box.text_frame
    p = tf.paragraphs[0]
    p.text = metrics_text
    p.font.size = Pt(TYPOGRAPHY["ppt_body_size"])  # 24pt

    # Sample info - readable label size
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(0.5))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"N = {summary.n_samples:,} | {summary.n_groups} demographic groups | Threshold: {summary.threshold:.0%}"
    p.font.size = Pt(TYPOGRAPHY["ppt_label_size"])  # 20pt


def _add_findings_slide(prs: Any, summary: AuditSummary) -> None:
    """Add key findings slide with publication-ready typography."""
    from pptx.util import Inches, Pt

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title - publication style (36pt)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Finding"
    p.font.size = Pt(TYPOGRAPHY["headline_size"])  # 36pt
    p.font.bold = True

    # Finding text - large readable body text
    disparity_pct = abs(summary.worst_disparity_value) * 100
    finding_text = (
        f"The largest disparity was found in {summary.worst_disparity_metric} "
        f"for the {summary.worst_disparity_group} group, "
        f"which differs by {disparity_pct:.1f}% from the reference group."
    )

    finding_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(3))
    tf = finding_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = finding_text
    p.font.size = Pt(TYPOGRAPHY["ppt_body_size"])  # 24pt - readable


def _add_recommendations_slide(prs: Any, summary: AuditSummary) -> None:
    """Add methodology slide with publication-ready typography."""
    from pptx.util import Inches, Pt

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title - publication style (36pt)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Methodology"
    p.font.size = Pt(TYPOGRAPHY["headline_size"])  # 36pt
    p.font.bold = True

    # Methodology information - readable body text
    methodology_text = [
        "Analysis Methodology",
        "",
        "Metrics computed per Van Calster B, Collins GS, Vickers AJ, et al.",
        '"Evaluation of performance measures in predictive artificial',
        'intelligence models to support medical decisions."',
        "Lancet Digit Health 2025.",
        "",
        "Results Summary:",
        f"  - Within threshold: {summary.pass_count} metrics",
        f"  - Near threshold: {summary.warn_count} metrics",
        f"  - Outside threshold: {summary.fail_count} metrics",
        "",
        "Healthcare organizations interpret these results based on",
        "clinical context, organizational values, and governance frameworks.",
    ]

    rec_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(4))
    tf = rec_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(methodology_text):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(TYPOGRAPHY["ppt_label_size"])  # 20pt - readable
        p.space_after = Pt(8)


def _add_exec_summary_chart_slide(prs: Any, results: "AuditResults") -> None:
    """Add executive summary figure slide."""
    _add_single_image_slide(prs, "Executive Summary", results.plot_executive_summary())


def _add_scorecard_chart_slide(prs: Any, results: "AuditResults") -> None:
    """Add go/no-go scorecard slide."""
    _add_single_image_slide(prs, "Go/No-Go Scorecard", results.plot_go_nogo_scorecard())


def _add_overall_charts_slide(prs: Any, results: "AuditResults") -> None:
    """Add 2x2 overall performance charts slide."""
    from faircareai.visualization.governance_dashboard import create_governance_overall_figures

    overall = create_governance_overall_figures(results)
    overall_figs = [fig for key, fig in overall.items() if key != "_explanations"]
    if len(overall_figs) >= 4:
        _add_grid_slide(prs, "Overall Performance", overall_figs[:4])
    elif overall_figs:
        _add_single_image_slide(prs, "Overall Performance", overall_figs[0])


def _add_subgroup_charts_slides(prs: Any, results: "AuditResults") -> None:
    """Add subgroup fairness slides (one per attribute)."""
    from faircareai.visualization.governance_dashboard import create_governance_subgroup_figures

    subgroup_figs = create_governance_subgroup_figures(results)
    for attr, fig_map in subgroup_figs.items():
        figs = list(fig_map.values())
        if len(figs) >= 4:
            _add_grid_slide(prs, f"Fairness by {attr}", figs[:4])
        elif figs:
            _add_single_image_slide(prs, f"Fairness by {attr}", figs[0])


def _add_vancalster_slide(prs: Any, results: "AuditResults") -> None:
    """Add Van Calster dashboard slide if possible."""
    if getattr(results, "_audit", None) is None:
        return
    audit = results._audit
    if not getattr(audit, "sensitive_attributes", None):
        return
    try:
        from faircareai.metrics.vancalster import compute_vancalster_metrics
        from faircareai.visualization.vancalster_plots import create_vancalster_dashboard

        vancalster = compute_vancalster_metrics(
            df=audit.df,
            y_prob_col=audit.pred_col,
            y_true_col=audit.target_col,
            group_col=audit.sensitive_attributes[0].column,
        )
        fig = create_vancalster_dashboard(vancalster)
        _add_single_image_slide(prs, "Van Calster Dashboard", fig)
    except Exception:
        return


def _add_single_image_slide(prs: Any, title: str, fig: Any) -> None:
    """Add a single chart slide."""
    from io import BytesIO

    from pptx.util import Inches

    from faircareai.reports.figure_exports import render_png_bytes

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _add_slide_title(slide, title)
    width_in = 12.3
    height_in = 5.9
    png = render_png_bytes(
        fig,
        scale=2,
        width=int(width_in * 150),
        height=int(height_in * 150),
    )
    image_stream = BytesIO(png)
    slide.shapes.add_picture(
        image_stream,
        Inches(0.5),
        Inches(1.1),
        width=Inches(width_in),
        height=Inches(height_in),
    )


def _add_grid_slide(prs: Any, title: str, figs: list[Any]) -> None:
    """Add a 2x2 grid of charts."""
    from io import BytesIO

    from pptx.util import Inches

    from faircareai.reports.figure_exports import render_png_bytes

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _add_slide_title(slide, title)

    left = Inches(0.5)
    top = Inches(1.1)
    gutter = Inches(0.3)
    cell_w_in = 6.0
    cell_h_in = 2.85
    cell_w = Inches(cell_w_in)
    cell_h = Inches(cell_h_in)

    positions = [
        (left, top),
        (left + cell_w + gutter, top),
        (left, top + cell_h + gutter),
        (left + cell_w + gutter, top + cell_h + gutter),
    ]

    for fig, (x, y) in zip(figs, positions, strict=False):
        png = render_png_bytes(
            fig,
            scale=2,
            width=int(cell_w_in * 150),
            height=int(cell_h_in * 150),
        )
        image_stream = BytesIO(png)
        slide.shapes.add_picture(image_stream, x, y, width=cell_w, height=cell_h)


# === Alias for PPTX generation ===
generate_pptx_report = generate_pptx_deck


# === GOVERNANCE PERSONA REPORT GENERATORS ===


def generate_governance_html_report(
    results: "AuditResults",
    output_path: str | Path,
    metric_config: "MetricDisplayConfig | None" = None,
    standalone: bool = True,
) -> Path:
    """Generate streamlined HTML report for governance committees.

    Creates a 3-5 page equivalent report with:
    - Executive Summary (traffic light status, pass/warn/fail counts)
    - Overall Performance (4 key figures)
    - Subgroup Performance (4 figures per attribute)
    - Recommendations and Sign-off

    Van Calster et al. (2025) Metric Display:
    -----------------------------------------
    Governance reports always show only RECOMMENDED metrics regardless of
    metric_config settings. The parameter is accepted for API consistency
    but OPTIONAL metrics are never shown in governance output.

    Args:
        results: AuditResults from FairCareAudit.run()
        output_path: Path for output HTML file
        metric_config: MetricDisplayConfig (ignored - governance shows RECOMMENDED only).

    Returns:
        Path to generated HTML file
    """
    output_path = _validate_output_path(Path(output_path))
    output_path.parent.mkdir(parents=True, exist_ok=True)

    html_content = _generate_governance_html(results)
    html_content = _inject_plotlyjs(html_content, standalone=standalone)

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(html_content)

    return output_path


def generate_governance_pdf_report(
    results: "AuditResults",
    output_path: str | Path,
    metric_config: "MetricDisplayConfig | None" = None,
) -> Path:
    """Generate streamlined PDF report for governance committees.

    Creates a 3-5 page report with key figures and plain language summaries.
    Uses Playwright to render interactive charts directly to PDF.

    Van Calster et al. (2025) Metric Display:
    -----------------------------------------
    Governance reports always show only RECOMMENDED metrics regardless of
    metric_config settings. The parameter is accepted for API consistency
    but OPTIONAL metrics are never shown in governance output.

    Args:
        results: AuditResults from FairCareAudit.run()
        output_path: Path for output PDF file
        metric_config: MetricDisplayConfig (ignored - governance shows RECOMMENDED only).

    Returns:
        Path to generated PDF file

    Raises:
        ImportError: If Playwright is not installed or chromium browser not available.
            Run: pip install playwright && playwright install chromium
    """
    try:
        from playwright.sync_api import sync_playwright  # noqa: F401
    except ImportError as err:
        raise ImportError(
            "Playwright is required for PDF generation. Install with: "
            "pip install 'faircareai[export]' && playwright install chromium"
        ) from err

    output_path = _validate_output_path(Path(output_path))
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Generate HTML with interactive charts
    html_content = _generate_governance_html(results)
    html_content = _inject_plotlyjs(html_content, standalone=True)

    # Use Playwright to render HTML to PDF (handles Jupyter/async context)
    _run_playwright_pdf_generation(html_content, output_path)

    return output_path


def _generate_governance_html(results: "AuditResults") -> str:
    """Generate streamlined HTML content for governance persona.

    Generates interactive Plotly charts that work in both HTML and PDF
    (when rendered with Playwright).

    Args:
        results: Audit results to render.
    """
    gov = results.governance_recommendation

    # Compute status from error/warning counts (don't rely on 'status' key)
    n_errors = gov.get("n_errors", gov.get("outside_threshold_count", 0))
    n_warnings = gov.get("n_warnings", gov.get("near_threshold_count", 0))

    if n_errors > 0:
        status = "REVIEW"
    elif n_warnings > 0:
        status = "CONDITIONAL"
    else:
        status = "READY"

    status_colors = {
        "READY": SEMANTIC_COLORS["pass"],
        "CONDITIONAL": SEMANTIC_COLORS["warn"],
        "REVIEW": SEMANTIC_COLORS["fail"],
    }
    status_color = status_colors.get(status, SEMANTIC_COLORS["fail"])
    # Use detection language - describe what was found, not what to do
    status_text = {
        "READY": "No Issues Detected",
        "CONDITIONAL": "Issues Near Threshold",
        "REVIEW": "Issues Exceeded Threshold",
    }
    report_generated_at = datetime.now().astimezone().isoformat(timespec="seconds")
    audit_run_at = results.run_timestamp or results.config.report_date or date.today().isoformat()

    # Generate interactive figures (work in both HTML and PDF via Playwright)
    try:
        overall_figures_html = _render_governance_overall_figures(results)
    except (ValueError, TypeError, KeyError) as e:
        logger.warning("Governance overall chart generation failed: %s", e)
        overall_figures_html = '<p class="chart-placeholder">Overall figures could not be generated.</p>'
    except FigureExportError as e:
        logger.warning("Chart export failed: %s", e)
        overall_figures_html = f'<p class="chart-placeholder">Chart export failed: {e.reason}</p>'
    except ImportError as e:
        logger.error("Visualization library missing: %s", e)
        overall_figures_html = '<p class="chart-placeholder">Install visualization dependencies: pip install \'faircareai[viz]\'</p>'
    except Exception:
        # Keep broad catch for truly unexpected errors
        logger.exception("Unexpected error in governance overall chart generation")
        overall_figures_html = (
            '<p class="chart-placeholder">Overall figures could not be generated.</p>'
        )

    try:
        subgroup_figures_html = _render_governance_subgroup_figures(results)
    except (ValueError, TypeError, KeyError) as e:
        logger.warning("Governance subgroup chart generation failed: %s", e)
        subgroup_figures_html = '<p class="chart-placeholder">Subgroup figures could not be generated.</p>'
    except FigureExportError as e:
        logger.warning("Chart export failed: %s", e)
        subgroup_figures_html = f'<p class="chart-placeholder">Chart export failed: {e.reason}</p>'
    except ImportError as e:
        logger.error("Visualization library missing: %s", e)
        subgroup_figures_html = '<p class="chart-placeholder">Install visualization dependencies: pip install \'faircareai[viz]\'</p>'
    except Exception:
        # Keep broad catch for truly unexpected errors
        logger.exception("Unexpected error in governance subgroup chart generation")
        subgroup_figures_html = '<p class="chart-placeholder">Subgroup figures could not be generated.</p>'

    # Plain language summary (use computed values from above)
    n_pass = gov.get("n_pass", gov.get("within_threshold_count", 0))

    # Generate plain language findings
    plain_findings = _generate_plain_language_findings(results)

    # Hero section data
    auroc_value = results.overall_performance.get("discrimination", {}).get("auroc", 0)
    hero_status = status_text.get(status, "Unknown")
    hero_color = status_color

    # Disparity callout data
    total_groups = gov.get("n_groups", 4)
    flag_count = n_errors
    worst_disparity = abs(gov.get("worst_disparity_value", 0))
    worst_metric = gov.get("worst_disparity_metric", "metric")
    worst_group = gov.get("worst_disparity_group", "group")

    # Primary fairness metric information
    from faircareai.core.config import FairnessMetric

    primary_metric = results.config.primary_fairness_metric
    metric_descriptions = {
        FairnessMetric.DEMOGRAPHIC_PARITY: ("Demographic Parity", "Equal selection rates across groups"),
        FairnessMetric.EQUALIZED_ODDS: ("Equalized Odds", "Equal TPR and FPR across groups"),
        FairnessMetric.EQUAL_OPPORTUNITY: ("Equal Opportunity", "Equal detection rates (TPR) across groups"),
        FairnessMetric.PREDICTIVE_PARITY: ("Predictive Parity", "Equal positive predictive values across groups"),
        FairnessMetric.CALIBRATION: ("Calibration", "Equal calibration accuracy across groups"),
    }
    metric_name, metric_desc = metric_descriptions.get(
        primary_metric, ("Not Specified", "No primary fairness metric was selected")
    )
    metric_justification = results.config.fairness_justification or "Not provided"

    audit_trail_html = _render_audit_trail_html(
        results,
        None,
        report_generated_at,
        title="Audit Trail",
    )

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>FairCareAI Governance Report: {results.config.model_name}</title>
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');

        :root {{
            --pass-color: {SEMANTIC_COLORS["pass"]};
            --warn-color: {SEMANTIC_COLORS["warn"]};
            --fail-color: {SEMANTIC_COLORS["fail"]};
            --bg-color: #ffffff;
            --text-color: #212529;
            --primary-color: #2c5282;
        }}

        * {{ box-sizing: border-box; }}

        body {{
            font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif;
            font-size: 16px;
            color: var(--text-color);
            background-color: var(--bg-color);
            line-height: 1.6;
            margin: 0;
            padding: 0;
        }}

        .container {{
            max-width: 900px;
            margin: 0 auto;
            padding: 40px 20px;
        }}

        h1 {{ font-size: 32px; font-weight: 700; margin-bottom: 8px; color: var(--primary-color); }}
        h2 {{ font-size: 24px; font-weight: 600; margin-top: 40px; border-bottom: 2px solid var(--primary-color); padding-bottom: 8px; }}
        h3 {{ font-size: 20px; font-weight: 600; margin-top: 24px; }}

        .header {{
            text-align: center;
            padding: 30px;
            border-bottom: 3px solid var(--primary-color);
            margin-bottom: 30px;
        }}

        .metadata {{ color: #666; font-size: 14px; margin-top: 8px; }}

        .status-badge {{
            display: inline-block;
            padding: 16px 32px;
            border-radius: 8px;
            font-weight: 700;
            font-size: 24px;
            color: white;
            background-color: {status_color};
            margin: 20px 0;
        }}

        .scorecard {{
            display: flex;
            justify-content: center;
            gap: 30px;
            margin: 30px 0;
        }}

        .scorecard-item {{
            text-align: center;
            padding: 20px 30px;
            border-radius: 8px;
            background: #f8f9fa;
            min-width: 120px;
        }}

        .scorecard-value {{
            font-size: 48px;
            font-weight: 700;
        }}

        .scorecard-label {{
            font-size: 14px;
            color: #666;
            text-transform: uppercase;
            letter-spacing: 0.5px;
        }}

        .pass {{ color: var(--pass-color); }}
        .warn {{ color: var(--warn-color); }}
        .fail {{ color: var(--fail-color); }}

        .section {{
            margin-bottom: 40px;
            page-break-inside: avoid;
        }}

        .findings-box {{
            background: #f8f9fa;
            padding: 24px;
            border-radius: 8px;
            border-left: 4px solid var(--primary-color);
            margin: 20px 0;
        }}

        .finding-item {{
            margin: 12px 0;
            padding: 8px 0;
            border-bottom: 1px solid #e2e8f0;
        }}

        .finding-item:last-child {{
            border-bottom: none;
        }}

        .figure-grid {{
            display: grid;
            grid-template-columns: repeat(2, 1fr);
            gap: 20px;
            margin: 20px 0;
        }}

        .figure-container {{
            background: #f8f9fa;
            padding: 16px;
            border-radius: 8px;
            text-align: center;
        }}

        .figure-title {{
            font-weight: 600;
            font-size: 14px;
            margin-bottom: 8px;
            color: var(--primary-color);
        }}

        .chart-placeholder {{
            background: #f0f0f0;
            padding: 60px 20px;
            text-align: center;
            border-radius: 8px;
            color: #666;
        }}

        .governance-block {{
            background: #f7fafc;
            border: 2px solid var(--primary-color);
            padding: 24px;
            border-radius: 8px;
            margin-top: 40px;
        }}



        .footer {{
            margin-top: 40px;
            padding: 20px;
            text-align: center;
            font-size: 14px;
            color: #666;
            border-top: 1px solid #e2e8f0;
        }}

        .audit-table {{
            width: 100%;
            border-collapse: collapse;
            font-size: 13px;
            margin-top: 12px;
        }}

        .audit-table th {{
            width: 32%;
            text-align: left;
            padding: 8px 10px;
            background: #f5f5f5;
            color: #555;
            font-weight: 600;
            border-bottom: 1px solid #e2e8f0;
            word-break: break-word;
            white-space: normal;
        }}

        .audit-table td {{
            padding: 8px 10px;
            border-bottom: 1px solid #e2e8f0;
            word-break: break-word;
            white-space: normal;
        }}

        .disclaimer {{
            font-size: 14px;
            color: #666;
            font-style: italic;
            background: #fffdf0;
            padding: 16px;
            border-radius: 6px;
            margin-top: 20px;
        }}

        /* Editorial-style Hero Section */
        .hero-section {{
            background: {hero_color};
            color: white;
            padding: 48px 40px;
            text-align: center;
            margin-bottom: 40px;
            border-radius: 8px;
        }}

        .hero-number {{
            font-size: 80px;
            font-weight: 700;
            line-height: 1;
            margin-bottom: 12px;
        }}

        .hero-title {{
            font-size: 28px;
            font-weight: 600;
            margin-bottom: 20px;
            opacity: 0.95;
        }}

        .hero-subtitle {{
            font-size: 18px;
            opacity: 0.9;
        }}

        /* Callout boxes for key statistics */
        .callout-box {{
            background: #fff3cd;
            border-left: 6px solid #ffc107;
            padding: 24px;
            margin: 30px 0;
            border-radius: 4px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.05);
        }}

        .callout-number {{
            font-size: 40px;
            font-weight: bold;
            color: #856404;
            margin-bottom: 8px;
        }}

        .callout-text {{
            font-size: 18px;
            color: #856404;
            margin-bottom: 12px;
        }}

        .callout-detail {{
            font-size: 14px;
            color: #666;
        }}

        /* Narrative section headlines */
        .narrative-headline {{
            font-size: 28px;
            font-weight: 600;
            color: var(--primary-color);
            margin: 40px 0 20px 0;
            border-bottom: 3px solid var(--primary-color);
            padding-bottom: 12px;
        }}

        @media print {{
            body {{ background: white; }}
            .container {{ max-width: 100%; }}
            .chart-grid, .figure-grid {{ page-break-inside: avoid; }}
            .section {{ page-break-inside: avoid; }}
            .hero-section {{ page-break-after: avoid; }}
        }}
    </style>
</head>
<body>
    <div class="container">
        <header class="header">
            <h1>Model Fairness Assessment</h1>
            <p style="font-size: 20px; color: #666;">Governance Committee Report</p>
            <p class="metadata">
                <strong>{results.config.model_name}</strong> v{results.config.model_version}<br>
                Report Date: {results.config.report_date or date.today().isoformat()}<br>
                Audit Run: {audit_run_at}<br>
                Report Generated: {report_generated_at}
            </p>
        </header>

        <!-- Information Banner -->
        <div style="background: #e7f3ff; border: 2px solid #0066cc; padding: 16px; border-radius: 8px; margin-bottom: 30px; text-align: center;">
            <strong style="font-size: 18px; color: #004080;">ℹ️ GOVERNANCE REVIEW MATERIALS</strong>
            <p style="margin: 8px 0 0 0; color: #004080;">
                This report provides statistical analysis and performance metrics for governance committee review.<br>
                Final deployment decisions are made by the health system governance team.
            </p>
        </div>

        <!-- Hero Section (editorial style) -->
        <div class="hero-section">
            <div class="hero-number">{auroc_value:.0%}</div>
            <div class="hero-title">Model Discrimination Score</div>
            <div class="hero-subtitle">{hero_status}</div>
        </div>

        <!-- Page 1: Executive Summary -->
        <section class="section">
            <h2>Executive Summary</h2>

            <div style="text-align: center;">
                <div class="status-badge">{status_text.get(status, status)}</div>
                <p style="font-size: 18px; margin-top: 16px; font-weight: 600;">
                    {_get_detection_summary(n_errors, n_warnings)}
                </p>
            </div>

            <div class="scorecard">
                <div class="scorecard-item">
                    <div class="scorecard-value pass">{n_pass}</div>
                    <div class="scorecard-label">Within Threshold</div>
                </div>
                <div class="scorecard-item">
                    <div class="scorecard-value warn">{n_warnings}</div>
                    <div class="scorecard-label">Near Threshold</div>
                </div>
                <div class="scorecard-item">
                    <div class="scorecard-value fail">{n_errors}</div>
                    <div class="scorecard-label">Exceeded Threshold</div>
                </div>
            </div>

            <div class="findings-box">
                <h3 style="margin-top: 0;">Key Findings</h3>
                {plain_findings}
            </div>

            <p class="disclaimer">
                This analysis follows the CHAI RAIC framework to provide performance and fairness metrics.
                The governance team will review these findings and make final deployment decisions through their established process.
            </p>
        </section>

        <!-- Page 2: Overall Performance -->
        <section class="section">
            <h2 class="narrative-headline">1. The Bottom Line: How Does the Model Perform?</h2>

            <div style="background: #f0f7ff; padding: 24px; margin: 20px 0; border-radius: 8px;">
                <p style="font-size: 18px; color: #333; margin: 0;">
                    <strong>In plain language:</strong> The model correctly ranks patients {auroc_value:.0%} of the time,
                    predicts risks accurately, and shows {hero_status.lower()}.
                </p>
            </div>

            <p style="color: #666; font-size: 14px; margin-bottom: 16px;">
                These 4 metrics tell the complete story - each gauge shows performance against established thresholds:
            </p>
            {overall_figures_html}
        </section>

        <!-- Pages 3-4: Subgroup Performance -->
        <section class="section">
            <h2 class="narrative-headline">2. Where Do Disparities Exist?</h2>

            <!-- Primary Fairness Metric Box -->
            <div style="background: #e8f4f8; border: 2px solid #0072B2; padding: 20px; margin-bottom: 24px; border-radius: 8px;">
                <h3 style="margin-top: 0; color: #0072B2; font-size: 18px;">Selected Fairness Metric: {metric_name}</h3>
                <p style="margin: 8px 0; color: #333;"><strong>Definition:</strong> {metric_desc}</p>
                <p style="margin: 8px 0 0 0; color: #666; font-size: 14px;"><strong>Justification:</strong> {metric_justification}</p>
            </div>

            <p style="color: #666; font-size: 16px; margin-bottom: 20px;">
                Performance varies across demographic groups. Charts corresponding to your selected metric are
                <span style="background: rgba(0, 114, 178, 0.1); padding: 2px 6px; border-radius: 3px;">highlighted in blue</span>.
            </p>

            <!-- Callout Box for Key Statistics -->
            <div class="callout-box">
                <div class="callout-number">{flag_count} of {total_groups}</div>
                <div class="callout-text">demographic groups flagged for review</div>
                <div class="callout-detail">
                    Largest disparity: {worst_disparity:.1%} in {worst_metric} for {worst_group}
                </div>
            </div>

            {subgroup_figures_html}

            <p style="color: #666; font-size: 14px; margin-top: 20px;">
                Bar charts show performance for each demographic group.
                Red bars indicate groups below threshold requiring attention.
            </p>
        </section>

        <!-- Page 5: Governance Decision -->
        <section class="section governance-block">
            <h2 class="narrative-headline">3. Your Decision: What Happens Next?</h2>

            <p><strong>Model:</strong> {results.config.model_name} v{results.config.model_version}</p>
            <p><strong>Intended Use:</strong> {results.config.intended_use or "Not specified"}</p>
            <p><strong>Intended Population:</strong> {results.config.intended_population or "Not specified"}</p>

            <div class="governance-note" style="background-color: #f8f9fa; border-left: 4px solid #0072B2; padding: 16px; margin-top: 20px;">
                <h3>Governance Process Note</h3>
                <p>
                    This package provides essential fairness and performance data to support the
                    <strong>Health System Governance Team's</strong> adjudication process.
                </p>
                <p>
                    The metrics and advisory status above are calculated outputs based on configured thresholds.
                    <strong>Final deployment decisions are made by the Governance Team outside of this technical package.</strong>
                </p>
                <p>
                    Please refer to your organization's standard operating procedures for voting, adjudication,
                    and final sign-off documentation.
                </p>
            </div>
        </section>

        {audit_trail_html}

        <footer class="footer">
            <p>{GOVERNANCE_DISCLAIMER_FULL}</p>
            <p style="font-size: 14px; margin-top: 12px;">
                <strong>Methodology:</strong> Van Calster B, Collins GS, Vickers AJ, et al.
                Evaluation of performance measures in predictive artificial intelligence models
                to support medical decisions: overview and guidance.
                <i>Lancet Digit Health</i> 2025;7(2):e100916.
                DOI: <a href="https://doi.org/10.1016/j.landig.2025.100916" target="_blank" style="color: #2c5282; text-decoration: none;">10.1016/j.landig.2025.100916</a>
            </p>
            <p>Generated by FairCareAI on {report_generated_at}</p>
        </footer>
    </div>
</body>
</html>"""

    return html


def _render_governance_overall_figures(results: "AuditResults") -> str:
    """Render the 4 overall performance figures for governance report.

    Generates interactive Plotly charts that work in both HTML and PDF
    (when rendered with Playwright).

    Args:
        results: AuditResults object
    """
    from faircareai.visualization.governance_dashboard import (
        create_governance_overall_figures,
    )

    try:
        figures = create_governance_overall_figures(results)

        # Remove _explanations dict (not a figure)
        figures.pop("_explanations", None)

        html_parts = ['<div class="figure-grid">']
        for title, fig in figures.items():
            if fig is not None and hasattr(fig, 'to_html'):
                # Render interactive Plotly chart
                fig_html = fig.to_html(full_html=False, include_plotlyjs=False)
                html_parts.append(f"""
                <div class="figure-container">
                    <div class="figure-title">{title}</div>
                    {fig_html}
                </div>
                """)
        html_parts.append("</div>")
        return "".join(html_parts)
    except (ValueError, TypeError, KeyError) as e:
        logger.warning("Governance overall figure rendering failed: %s", e)
        return f'<p class="chart-placeholder">Overall figures could not be generated: {html.escape(str(e))}</p>'
    except FigureExportError as e:
        logger.warning("Chart export failed: %s", e)
        return f'<p class="chart-placeholder">Chart export failed: {e.reason}</p>'
    except ImportError as e:
        logger.error("Visualization library missing: %s", e)
        return '<p class="chart-placeholder">Install visualization dependencies: pip install \'faircareai[viz]\'</p>'
    except Exception as e:
        # Keep broad catch for truly unexpected errors
        logger.exception("Unexpected error rendering governance overall figures")
        return f'<p class="chart-placeholder">Overall figures could not be generated: {html.escape(str(e))}</p>'


def _render_governance_subgroup_figures(results: "AuditResults") -> str:
    """Render the subgroup performance figures for governance report.

    Generates interactive Plotly charts that work in both HTML and PDF
    (when rendered with Playwright).

    Args:
        results: Audit results to visualize.
    """
    from faircareai.visualization.governance_dashboard import (
        create_governance_subgroup_figures,
    )

    try:
        # Get figures for each sensitive attribute
        all_figures = create_governance_subgroup_figures(results)

        html_parts = []
        for attr_name, figures in all_figures.items():
            html_parts.append(f"<h3>{attr_name.replace('_', ' ').title()}</h3>")
            html_parts.append('<div class="figure-grid">')
            for title, fig in figures.items():
                if fig is not None:
                    # Render interactive Plotly chart
                    fig_html = fig.to_html(full_html=False, include_plotlyjs=False)
                    html_parts.append(f"""
                    <div class="figure-container">
                        <div class="figure-title">{title}</div>
                        {fig_html}
                    </div>
                    """)
            html_parts.append("</div>")

        return (
            "".join(html_parts)
            if html_parts
            else '<p class="chart-placeholder">No subgroup figures available.</p>'
        )
    except (ValueError, TypeError, KeyError) as e:
        logger.warning("Governance subgroup figure rendering failed: %s", e)
        return f'<p class="chart-placeholder">Subgroup figures could not be generated: {html.escape(str(e))}</p>'
    except FigureExportError as e:
        logger.warning("Chart export failed: %s", e)
        return f'<p class="chart-placeholder">Chart export failed: {e.reason}</p>'
    except ImportError as e:
        logger.error("Visualization library missing: %s", e)
        return '<p class="chart-placeholder">Install visualization dependencies: pip install \'faircareai[viz]\'</p>'
    except Exception as e:
        # Keep broad catch for truly unexpected errors
        logger.exception("Unexpected error rendering governance subgroup figures")
        return f'<p class="chart-placeholder">Subgroup figures could not be generated: {html.escape(str(e))}</p>'


def _generate_plain_language_findings(results: "AuditResults") -> str:
    """Generate plain language findings for governance report."""
    perf = results.overall_performance
    disc = perf.get("discrimination", {})
    cal = perf.get("calibration", {})
    gov = results.governance_recommendation

    findings = []

    # Model discrimination
    auroc = disc.get("auroc", 0)
    if auroc >= 0.8:
        findings.append(
            "The model demonstrates <strong>strong</strong> ability to distinguish between outcomes."
        )
    elif auroc >= 0.7:
        findings.append(
            "The model demonstrates <strong>acceptable</strong> ability to distinguish between outcomes."
        )
    else:
        findings.append(
            "The model's ability to distinguish between outcomes <strong>requires review</strong>."
        )

    # Calibration
    brier = cal.get("brier_score", 1)
    slope = cal.get("calibration_slope", 0)
    if brier < 0.15 and 0.8 <= slope <= 1.2:
        findings.append(
            "Model predictions are <strong>well-calibrated</strong> with actual outcomes."
        )
    elif brier < 0.25:
        findings.append("Model calibration is <strong>acceptable</strong> but could be improved.")
    else:
        findings.append("Model calibration <strong>requires attention</strong>.")

    # Fairness - use detection language, not recommendation language
    n_errors = gov.get("n_errors", 0)
    n_warnings = gov.get("n_warnings", 0)
    if n_errors == 0 and n_warnings == 0:
        findings.append("No fairness disparities detected across demographic groups.")
    elif n_errors == 0:
        findings.append(
            f"<strong>{n_warnings} fairness metric(s)</strong> detected near threshold."
        )
    else:
        findings.append(
            f"<strong>{n_errors} fairness metric(s)</strong> detected exceeding threshold."
        )

    # Cohort size
    n_total = results.descriptive_stats.get("cohort_overview", {}).get("n_total", 0)
    if n_total >= 1000:
        findings.append(f"Analysis based on adequate sample size (N = {n_total:,}).")
    elif n_total >= 100:
        findings.append(
            f"Sample size (N = {n_total:,}) is limited; results should be interpreted with caution."
        )
    else:
        findings.append(
            f"<strong>Warning:</strong> Small sample size (N = {n_total:,}) limits reliability of results."
        )

    return ''.join(f'<div class="finding-item">{finding}</div>' for finding in findings)


def _get_detection_summary(n_errors: int, n_warnings: int) -> str:
    """Generate detection language summary for governance report.

    Uses "X metrics exceeded threshold" language instead of pass/fail.

    Args:
        n_errors: Number of metrics that exceeded threshold.
        n_warnings: Number of metrics near threshold.

    Returns:
        Plain language detection summary.
    """
    if n_errors == 0 and n_warnings == 0:
        return "No fairness issues detected"
    elif n_errors == 0:
        return f"{n_warnings} metric(s) near threshold"
    elif n_warnings == 0:
        return f"{n_errors} metric(s) exceeded threshold"
    else:
        return f"{n_errors} metric(s) exceeded threshold, {n_warnings} near threshold"


def _get_governance_print_css() -> str:
    """Get CSS for governance PDF output (3-5 pages)."""
    return """
    @page {
        size: letter;
        margin: 0.75in;

        @top-right {
            content: "Governance Report";
            font-size: 12pt;
            color: #666;
        }

        @bottom-center {
            content: counter(page) " of " counter(pages);
            font-size: 12pt;
            color: #666;
        }
    }

    @page :first {
        @top-right { content: none; }
    }

    body {
        font-size: 12pt;
    }

    h1 { font-size: 24pt; }
    h2 { font-size: 18pt; page-break-after: avoid; }
    h3 { font-size: 14pt; }

    .section {
        page-break-inside: avoid;
    }

    .chart-grid, .figure-grid {
        page-break-inside: avoid;
    }

    .governance-block {
        page-break-before: always;
    }
    """
